"""基于真实 PPT 模板（deck）的槽位填字渲染器。

加载 templates/PPT_template/*.pptx 成品模板，按 templates/ppt_decks/deck_slots.json
的角色+位置锚点，把生成内容填入对应文本形状，保留模板的视觉设计（装饰/图表/配图）。
"""
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

SLOTS_PATH = Path(__file__).resolve().parents[3] / "templates" / "ppt_decks" / "deck_slots.json"
DECK_TEMPLATES_DIR = Path(__file__).resolve().parents[3] / "templates" / "PPT_template"
EMU_PER_INCH = 914400


def _load_slots_data() -> dict:
    return json.loads(SLOTS_PATH.read_text(encoding="utf-8"))


def _load_slots(template_id: str) -> dict:
    """按模板 id 读取该模板的槽位几何（每套模板版式独立）。"""
    data = _load_slots_data()
    templates = data.get("templates") or {}
    if template_id not in templates:
        raise KeyError(f"deck_slots.json 缺少模板 {template_id} 的槽位几何")
    return templates[template_id]


def role_order() -> list[str]:
    return _load_slots_data()["role_order"]


def slot_counts(template_id: str) -> dict[str, int]:
    """返回模板每个角色的内容槽位数，供 make_deck 按模板整形内容。"""
    roles = _load_slots(template_id)
    return {role: len(spec["content"]) for role, spec in roles.items()}


# deck 角色 → 现有 page_type 映射（不新增 schema 枚举）。
ROLE_PAGE_TYPE: dict[str, str] = {
    "cover": "cover", "intro": "scenario", "objectives": "objectives",
    "knowledge_map": "concept", "knowledge_intro": "concept",
    "core_1": "concept", "core_2": "concept", "core_3": "concept", "core_4": "concept",
    "case_study": "case", "discussion": "question", "summary": "summary",
    "assessment": "exercise", "assignment": "homework", "end": "summary",
}
PAGE_PURPOSE: dict[str, str] = {
    "cover": "建立课程主题并唤起学习期待",
    "scenario": "用真实问题激活学生既有经验",
    "objectives": "明确本课可观察的学习成果",
    "concept": "建立核心概念与关键关系",
    "case": "通过案例示范应用过程",
    "question": "引导互动讨论与判断",
    "exercise": "即时测验收集学习证据",
    "summary": "总结本课要点并预告联系",
    "homework": "布置课后应用任务",
}
PAGE_LAYOUT: dict[str, str] = {
    "cover": "cover", "scenario": "question", "objectives": "bullet",
    "concept": "split", "case": "split", "question": "question",
    "exercise": "exercise", "summary": "summary", "homework": "bullet",
}
PAGE_VISUAL: dict[str, str] = {
    "cover": "纯白背景左侧主色竖栏，课程主题大标题与副标题留白排版",
    "scenario": "上方大面积留白作为问题区，下方虚线框提示学生写下初步判断",
    "objectives": "编号列表纵向排列，编号使用主题色圆形徽章",
    "concept": "左侧概念框图，右侧箭头图表示关键关系，底部保留留白",
    "case": "横向步骤流程线，每步配编号与短句",
    "question": "上方问题区，下方作答提示条",
    "exercise": "左题目右作答分栏，底部自我检查提示条",
    "summary": "三条结论短句，下方时间轴示意环节推进",
    "homework": "任务清单卡片，逐条核对",
}
ROLE_WEIGHT: dict[str, float] = {
    "cover": 0.04, "intro": 0.10, "objectives": 0.08, "knowledge_map": 0.08,
    "knowledge_intro": 0.10, "core_1": 0.07, "core_2": 0.07, "core_3": 0.07, "core_4": 0.07,
    "case_study": 0.10, "discussion": 0.07, "summary": 0.06, "assessment": 0.06,
    "assignment": 0.03, "end": 0.02,
}


def deck_structure(template_id: str) -> dict:
    """返回模板的 deck 角色结构，供 LLM 生成内容时对齐真实模板的页序与槽位。

    每项含角色名、映射 page_type、页面目的、版式、视觉建议与内容槽位数。
    slides[i]（i 从 0 起）对应模板第 i+1 页；LLM 按此结构生成恰好
    ``page_count`` 页、页序不可改变的内容。
    """
    counts = slot_counts(template_id)
    roles = []
    for index, role in enumerate(role_order(), 1):
        page_type = ROLE_PAGE_TYPE.get(role, "concept")
        roles.append({
            "index": index,
            "role": role,
            "page_type": page_type,
            "purpose": PAGE_PURPOSE.get(page_type, "讲解要点"),
            "layout": PAGE_LAYOUT.get(page_type, "bullet"),
            "visual_suggestion": PAGE_VISUAL.get(page_type, "要点列表"),
            "slot_count": counts.get(role, 0),
        })
    return {
        "template_id": template_id,
        "page_count": len(roles),
        "roles": roles,
    }


def _find_shape(slide, anchor: dict, tol_inches: float):
    """返回最接近锚点 (x,y) 英寸坐标的文本形状，无则 None。

    只匹配文本框/占位符，跳过色块、圆、图片等装饰形状，避免把内容填进装饰元素。
    """
    tx = anchor["x"] * EMU_PER_INCH
    ty = anchor["y"] * EMU_PER_INCH
    best = None
    best_dist = None
    for shape in slide.shapes:
        if shape.left is None or shape.top is None or not shape.has_text_frame:
            continue
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX and not shape.is_placeholder:
            continue
        dist = abs(shape.left - tx) + abs(shape.top - ty)
        if best_dist is None or dist < best_dist:
            best, best_dist = shape, dist
    if best is not None and best_dist <= tol_inches * EMU_PER_INCH * 2:
        return best
    return None


def _fit_font_size(shape, text: str, base_size: float) -> float:
    """按文本框尺寸估算并收缩字号，避免长文本溢出。"""
    width_in = (shape.width or Emu(int(6 * EMU_PER_INCH))) / EMU_PER_INCH
    height_in = (shape.height or Emu(int(1 * EMU_PER_INCH))) / EMU_PER_INCH
    size = float(base_size)
    while size > 8:
        # 中文字符宽度约等于字号；按整串字符估算（拉丁字符偏保守）
        char_w = size / 72.0 * 0.98
        chars_per_line = max(1, int(width_in / char_w))
        lines = 0
        for segment in text.split("\n"):
            lines += max(1, math.ceil(len(segment) / chars_per_line))
        needed_height = lines * size / 72.0 * 1.28
        if needed_height <= height_in * 0.98:
            break
        size -= 1
    return size


def _set_shape_text(shape, value: str):
    """写入文本并保留首段首个 run 的字体格式，清除多余 run/段落，并收缩字号防止溢出。"""
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    base_size = 16.0
    if paragraph.runs:
        first_run = paragraph.runs[0]
        if first_run.font.size:
            base_size = float(first_run.font.size.pt)
        first_run.text = str(value)
        for run in paragraph.runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        paragraph.text = str(value)
    for extra in frame.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    frame.word_wrap = True
    size = _fit_font_size(shape, str(value), base_size)
    if paragraph.runs:
        paragraph.runs[0].font.size = Pt(size)
    else:
        paragraph.font.size = Pt(size)


def render_deck(template_path: Path, slides: list[dict], output: Path, template_id: str) -> Path:
    """把 slides 内容填入模板 deck，输出新 .pptx。

    slides[i] 对应模板第 i+1 页（i 与 role_order 对齐）；每页填 title 槽 +
    content 槽（顺序取自 slides[i]["body"]）。槽位几何按 template_id 独立读取，
    因此每套模板的版式（封面/分栏/卡片/装饰）各不相同。
    """
    roles = _load_slots(template_id)
    order = role_order()
    tol = 0.35
    presentation = Presentation(str(template_path))
    for index, slide in enumerate(slides[: len(presentation.slides)]):
        if index >= len(order):
            break
        role = order[index]
        spec = roles.get(role)
        if not spec:
            continue
        template_slide = presentation.slides[index]
        title_shape = _find_shape(template_slide, spec["title"], tol)
        if title_shape is not None:
            _set_shape_text(title_shape, str(slide.get("title") or ""))
        body = [str(value) for value in (slide.get("body") or [])]
        for anchor_index, anchor in enumerate(spec["content"]):
            if anchor_index >= len(body):
                break
            shape = _find_shape(template_slide, anchor, tol)
            if shape is not None:
                _set_shape_text(shape, body[anchor_index])
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    return output


def deck_template_path(template_id: str) -> Path:
    """把模板 id 解析为 .pptx 路径。

    兼容全量模板 id（lessonforge_deck_academic）与短名（Academic_Template / academic）。
    """
    name = template_id
    if name.startswith("lessonforge_deck_"):
        name = name[len("lessonforge_deck_"):]
    candidates = [
        DECK_TEMPLATES_DIR / f"{name}.pptx",
        DECK_TEMPLATES_DIR / f"{name}_Template.pptx",
        DECK_TEMPLATES_DIR / f"{name.capitalize()}_Template.pptx",
    ]
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    # 按文件前缀宽松匹配（忽略大小写/下划线）
    for candidate in DECK_TEMPLATES_DIR.glob("*.pptx"):
        if candidate.stem.lower().replace("_", "") == name.lower().replace("_", ""):
            return candidate
    raise FileNotFoundError(f"未找到 deck 模板：{template_id}")
