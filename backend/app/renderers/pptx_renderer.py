from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.services.ppt_template_service import resolve_ppt_template


SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


def _rgb(value: str) -> RGBColor:
    value = value.removeprefix("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _shape(slide, shape_type, x, y, width, height, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type,
        Inches(x), Inches(y), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def _textbox(
    slide,
    text,
    x,
    y,
    width,
    height,
    *,
    font,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    margin=0,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _decorate(slide, template: dict, index: int, total: int):
    colors = {key: _rgb(value) for key, value in template["palette"].items()}
    fonts = template["typography"]
    mode = template["composition"]
    if mode == "swiss_rail":
        _shape(slide, MSO_SHAPE.RECTANGLE, .55, .55, .08, 6.38, colors["primary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.8, .55, 10.85, .02, colors["secondary"])
    elif mode == "nordic_field":
        _shape(slide, MSO_SHAPE.OVAL, 11.15, .18, 1.85, 1.85, colors["secondary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, .7, 6.82, 3.1, .08, colors["primary"])
    elif mode == "academic_offset":
        _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 2.0, 7.5, colors["primary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 2.0, 0, .12, 7.5, colors["secondary"])
    elif mode == "editorial_margin":
        _shape(slide, MSO_SHAPE.RECTANGLE, .72, .65, .03, 6.15, colors["primary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.0, 1.32, 11.55, .02, colors["secondary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, 1.0, 6.62, 11.55, .02, colors["secondary"])
    elif mode == "science_signal":
        _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, .11, colors["primary"])
        _shape(slide, MSO_SHAPE.RECTANGLE, .72, .68, .06, 5.95, colors["primary"])
        for offset in range(5):
            _shape(slide, MSO_SHAPE.RECTANGLE, 9.1 + offset * .7, 6.85, .42, .035, colors["secondary"])
    elif mode == "primary_blocks":
        _shape(slide, MSO_SHAPE.OVAL, .08, .08, 1.18, 1.18, colors["secondary"])
        _shape(slide, MSO_SHAPE.OVAL, 12.18, 6.38, .9, .9, colors["primary"])

    folio_x = .38 if mode == "academic_offset" else .44 if mode == "primary_blocks" else .82
    folio_color = colors["on_primary"] if mode in {"academic_offset", "primary_blocks"} else colors["primary"]
    _textbox(
        slide, f"{index:02d}", folio_x, .62, .9, .42,
        font=fonts["latin"], size=16, color=folio_color, bold=True,
    )
    footer_x = 11.1 if mode == "primary_blocks" else 11.55
    _textbox(
        slide, f"{index} / {total}", footer_x, 6.88, .8, .25,
        font=fonts["latin"], size=9, color=colors["muted"], align=PP_ALIGN.RIGHT,
    )
    return colors, fonts


def _render_cover(slide, item: dict, template: dict, colors: dict, fonts: dict):
    mode = template["composition"]
    x = 2.55 if mode == "academic_offset" else 1.35
    width = 9.7 if mode == "academic_offset" else 10.6
    accent_color = colors["on_primary"] if mode == "academic_offset" else colors["primary"]
    if mode == "academic_offset":
        _textbox(
            slide, "LESSONFORGE", .35, 5.95, 1.3, .35,
            font=fonts["latin"], size=10, color=accent_color, bold=True,
        )
    _textbox(
        slide, item.get("title", ""), x, 1.7, width, 1.55,
        font=fonts["heading"], size=42, color=colors["text"], bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    body = item.get("body") or []
    subtitle = " · ".join(str(value) for value in body[:2])
    if subtitle:
        _textbox(
            slide, subtitle, x, 3.42, width, .62,
            font=fonts["body"], size=22, color=colors["muted"],
        )
    purpose = item.get("purpose")
    if purpose:
        _textbox(
            slide, purpose, x, 5.62, 5.6, .42,
            font=fonts["body"], size=16, color=colors["primary"], bold=True,
        )


def _render_process(slide, body: list[str], colors: dict, fonts: dict, x: float, width: float):
    count = max(1, min(len(body), 4))
    gap = .22
    block_width = (width - gap * (count - 1)) / count
    for item_index, text in enumerate(body[:4], 1):
        left = x + (item_index - 1) * (block_width + gap)
        _textbox(
            slide, f"{item_index:02d}", left, 2.12, block_width, .55,
            font=fonts["latin"], size=26, color=colors["primary"], bold=True,
        )
        _shape(slide, MSO_SHAPE.RECTANGLE, left, 2.85, block_width, .04, colors["secondary"])
        _textbox(
            slide, text, left, 3.18, block_width, 1.65,
            font=fonts["body"], size=20, color=colors["text"], bold=True,
        )


def _render_standard(slide, item: dict, template: dict, colors: dict, fonts: dict):
    mode = template["composition"]
    x = 2.55 if mode == "academic_offset" else 1.35
    width = 9.7 if mode == "academic_offset" else 10.6
    _textbox(
        slide, item.get("title", ""), x, .88, width, .82,
        font=fonts["heading"], size=34, color=colors["text"], bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    body = [str(value) for value in (item.get("body") or [])]
    page_type = item.get("page_type", "concept")
    layout = item.get("layout", "")
    if page_type == "process" or layout == "steps":
        _render_process(slide, body, colors, fonts, x, width)
    elif page_type in {"question", "exercise"}:
        _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.0, width, 3.8, colors["surface"], colors["secondary"], radius=True)
        for body_index, text in enumerate(body[:4], 1):
            _textbox(
                slide, f"{body_index}", x + .35, 2.35 + (body_index - 1) * .78, .42, .42,
                font=fonts["latin"], size=15, color=colors["on_primary"], bold=True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
            ).fill.solid()
            badge = slide.shapes[-1]
            badge.fill.fore_color.rgb = colors["primary"]
            _textbox(
                slide, text, x + .9, 2.3 + (body_index - 1) * .78, width - 1.25, .56,
                font=fonts["body"], size=20, color=colors["text"],
            )
    elif page_type == "comparison" or layout == "split":
        split = max(1, (len(body) + 1) // 2)
        for column, values in enumerate((body[:split], body[split:])):
            left = x + column * (width / 2 + .12)
            column_width = width / 2 - .18
            _shape(slide, MSO_SHAPE.RECTANGLE, left, 2.08, column_width, .06, colors["primary"] if column == 0 else colors["secondary"])
            for body_index, text in enumerate(values[:4]):
                _textbox(
                    slide, text, left, 2.55 + body_index * .78, column_width, .62,
                    font=fonts["body"], size=20, color=colors["text"],
                )
    else:
        body_box = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(width), Inches(3.85))
        frame = body_box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(.05)
        for body_index, text in enumerate(body[:6]):
            paragraph = frame.paragraphs[0] if body_index == 0 else frame.add_paragraph()
            paragraph.text = text
            paragraph.level = 0
            paragraph.font.name = fonts["body"]
            paragraph.font.size = Pt(21)
            paragraph.font.color.rgb = colors["text"]
            paragraph.space_after = Pt(15)
            clean_text = text.lstrip("•-*123456789. ").strip()
            paragraph.text = f"•  {clean_text}"

    suggestion = item.get("visual_suggestion")
    if suggestion:
        _shape(slide, MSO_SHAPE.RECTANGLE, x, 6.16, width, .5, colors["surface"])
        _textbox(
            slide, suggestion, x + .16, 6.25, width - .32, .28,
            font=fonts["body"], size=10, color=colors["muted"],
        )


def render_pptx(title: str, content: dict, output: Path) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)
    template = resolve_ppt_template(content.get("theme"))
    slides = content.get("slides", [])
    for index, item in enumerate(slides, 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = _rgb(template["palette"]["background"])
        colors, fonts = _decorate(slide, template, index, len(slides))
        if item.get("page_type") == "cover" or item.get("layout") in {"title", "cover"}:
            _render_cover(slide, item, template, colors, fonts)
        else:
            _render_standard(slide, item, template, colors, fonts)
        slide.notes_slide.notes_text_frame.text = item.get("speaker_notes", "")
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    Presentation(output)
    return output
