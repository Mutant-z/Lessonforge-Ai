"""动态 PPT 演示构建器（PPT Editing Agent 的工具层）。

模板被理解为「设计语言来源」（调色板 / 字体 / 装饰元素），不是固定填充容器。
Agent 通过编辑工具（create_slide / add_textbox / add_shape / add_image / add_chart /
move / resize / delete / set_style / set_background / add_notes）逐步构建每页元素，
render() 按元素几何 + 模板装饰程序化生成真实 PPTX。

坐标系统：英寸，画布 13.333 × 7.5（与 pptx_renderer.SLIDE_WIDTH/HEIGHT 一致）。
"""
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.services.ppt_template_service import resolve_ppt_template
from app.agent.slide_rendering import infer_render_mode, semantic_body_texts
from app.renderers.deck_renderer import deck_structure

SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
EMU_PER_INCH = 914400
FIT_FLOOR_PT = 8.0


# 每套模板内容页/封面/尾页的装饰元素（源自 scripts/build_generic_decks.py 的 shell/cover/end）
# 装饰规格：(kind, x, y, w, h, color_key)，kind ∈ rect|oval
TEMPLATE_DECOR: dict[str, dict[str, Any]] = {
    "lessonforge_deck_academic": {
        "shell": [("rect", 0, 0, 0.16, 7.5, "primary")],
        "cover": [("rect", 0, 0, 2.1, 7.5, "primary"), ("rect", 2.1, 0, 0.06, 7.5, "secondary")],
        "end": [("oval", -1.5, 4.5, 4.8, 4.8, "secondary")],
    },
    "lessonforge_deck_ai_future": {
        "shell": [("rect", 0, 0, 13.333, 0.12, "primary"), ("oval", 12.6, 0.35, 0.5, 0.5, "secondary")],
        "cover": [("rect", 0, 0, 13.333, 0.1, "secondary"), ("rect", 1.0, 1.9, 11.3, 0.02, "secondary"), ("rect", 0.6, 6.4, 6.0, 0.02, "secondary")],
        "end": [("rect", 0, 0, 13.333, 0.1, "secondary"), ("oval", 11.4, -1.4, 3.6, 3.6, "secondary")],
    },
    "lessonforge_deck_business": {
        "shell": [("rect", 0, 0, 13.333, 0.1, "primary"), ("rect", 0.65, 6.85, 12.0, 0.015, "secondary")],
        "cover": [("rect", 0, 0, 13.333, 0.35, "primary"), ("rect", 1.0, 3.95, 2.2, 0.09, "secondary")],
        "end": [("rect", 0, 0, 13.333, 0.35, "secondary")],
    },
    "lessonforge_deck_cartoon": {
        "shell": [("oval", 12.3, 0.4, 0.7, 0.7, "secondary")],
        "cover": [("oval", 11.0, 0.4, 1.6, 1.6, "secondary"), ("oval", 0.7, 0.5, 1.1, 1.1, "secondary"), ("rect", 1.0, 3.6, 3.0, 0.12, "secondary")],
        "end": [("oval", 11.2, -1.2, 3.2, 3.2, "secondary"), ("oval", -0.8, 5.6, 2.6, 2.6, "secondary")],
    },
    "lessonforge_deck_chinese_culture": {
        "shell": [("rect", 0, 0, 0.18, 7.5, "primary")],
        "cover": [("rect", 0, 0, 0.3, 7.5, "primary"), ("rect", 0.3, 0, 0.08, 7.5, "secondary"), ("rect", 1.6, 4.0, 2.4, 0.09, "primary")],
        "end": [("rect", 0, 0, 0.3, 7.5, "secondary"), ("oval", 11.3, -1.2, 3.2, 3.2, "secondary")],
    },
    "lessonforge_deck_smart_ai": {
        "shell": [("rect", 0, 0, 1.7, 7.5, "primary"), ("rect", 1.7, 0, 0.07, 7.5, "secondary")],
        "cover": [("rect", 0, 0, 2.6, 7.5, "secondary")],
        "end": [("rect", 0, 0, 2.6, 7.5, "secondary"), ("oval", 11.4, -1.2, 3.4, 3.4, "secondary")],
    },
}


def _rgb(value: str | None, fallback: str = "#000000") -> RGBColor:
    value = (value or fallback).removeprefix("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _color(palette: dict[str, str], key: str) -> RGBColor:
    value = palette.get(key, "#000000")
    try:
        return _rgb(value)
    except ValueError:
        return RGBColor(0, 0, 0)


def design_system_for(template: dict[str, Any]) -> dict[str, Any]:
    """把 catalog 模板元数据扩展为可渲染的设计系统（含装饰几何）。

    deck 模板（composition=="deck"）额外携带 deck_structure：15 页角色顺序、
    每页映射 page_type/目的/版式/视觉建议与内容槽位数，供 LLM 生成内容时
    对齐真实模板的页序与槽位，也让所有 Agent 能感知模板的页面结构。
    """
    template_id = template["id"]
    decor = TEMPLATE_DECOR.get(template_id, TEMPLATE_DECOR["lessonforge_deck_academic"])
    design = {
        "id": template_id,
        "palette": template.get("palette", {}),
        "typography": template.get("typography", {}),
        "composition": template.get("composition", "deck"),
        "decoration": decor,
        "safe_margin": {"x": 0.6, "y": 0.5, "bottom": 0.7},
        "canvas": {"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT},
    }
    if design["composition"] == "deck":
        design["deck_structure"] = deck_structure(template_id)
    return design


class PresentationBuilder:
    """可编程构建 / 编辑演示文稿的模型层（不直接操作 XML）。"""

    def __init__(self, template_id: str | None = None):
        self.template = resolve_ppt_template(template_id)
        self.design_system = design_system_for(self.template)
        self.slides: list[dict[str, Any]] = []
        self._element_seq = 0

    # ---------- 幻灯片与元素编辑 ----------
    def create_slide(self, page_type: str = "concept", title: str = "", layout: str = "bullet", purpose: str = "") -> str:
        slide_id = f"S{len(self.slides) + 1:02d}"
        self.slides.append({
            "id": slide_id, "page_type": page_type, "title": title, "purpose": purpose,
            "body": [], "layout": layout, "visual_suggestion": "",
            "speaker_notes": "", "duration_seconds": 0, "blocks": [],
            "script_segment_ids": [],
            "background": None, "elements": [],
            "render_mode": "semantic",
        })
        return slide_id

    def get_slide(self, slide_id: str) -> dict[str, Any]:
        for slide in self.slides:
            if slide["id"] == slide_id:
                return slide
        raise KeyError(f"幻灯片不存在：{slide_id}")

    def _add_element(self, slide_id: str, kind: str, x: float, y: float, w: float, h: float, **fields) -> str:
        slide = self.get_slide(slide_id)
        self._element_seq += 1
        element_id = f"E{self._element_seq:02d}"
        element = {
            "id": element_id, "kind": kind, "x": x, "y": y, "w": w, "h": h,
            "z": fields.pop("z", self._element_seq),
            "style": fields.pop("style", {}),
            **fields,
        }
        slide["elements"].append(element)
        return element_id

    def set_slide_title(self, slide_id: str, title: str) -> str:
        slide = self.get_slide(slide_id)
        slide["title"] = title
        return slide_id

    def add_textbox(self, slide_id: str, text: str, x: float, y: float, w: float, h: float,
                    style: dict[str, Any] | None = None, role: str = "", content_ref: str = "") -> str:
        # A caller that starts placing textboxes is defining the complete
        # editable geometry layer. Media-only additions use hybrid instead.
        self.get_slide(slide_id)["render_mode"] = "absolute"
        return self._add_element(
            slide_id, "textbox", x, y, w, h, text=text, style=style or {}, role=role,
            content_ref=content_ref,
        )

    def add_shape(self, slide_id: str, shape_type: str = "rect", x: float = 0, y: float = 0,
                  w: float = 1, h: float = 1, fill: str | None = None, line: str | None = None,
                  radius: bool = False, role: str = "") -> str:
        return self._add_element(
            slide_id, "shape", x, y, w, h,
            shape_type=shape_type, fill=fill, line=line, radius=radius, role=role,
        )

    def add_image(self, slide_id: str, file_path: str, x: float, y: float, w: float, h: float,
                  role: str = "image", asset_id: str = "", provider: str = "", degraded: bool = False,
                  visual_slot: str = "primary_visual") -> str:
        slide = self.get_slide(slide_id)
        if infer_render_mode(slide) == "semantic":
            slide["render_mode"] = "hybrid"
        return self._add_element(
            slide_id, "image", x, y, w, h,
            asset_path=file_path, asset_id=asset_id, provider=provider, degraded=degraded, role=role,
            visual_slot=visual_slot,
        )

    def add_chart(self, slide_id: str, chart_type: str, data: dict[str, Any], x: float, y: float,
                  w: float, h: float, role: str = "chart") -> str:
        slide = self.get_slide(slide_id)
        if infer_render_mode(slide) == "semantic":
            slide["render_mode"] = "hybrid"
        return self._add_element(slide_id, "chart", x, y, w, h, chart={"type": chart_type, "data": data}, role=role)

    def move_element(self, slide_id: str, element_id: str, x: float, y: float) -> str:
        element = self._find_element(slide_id, element_id)
        element["x"], element["y"] = x, y
        return element_id

    def resize_element(self, slide_id: str, element_id: str, w: float, h: float) -> str:
        element = self._find_element(slide_id, element_id)
        element["w"], element["h"] = w, h
        return element_id

    def delete_element(self, slide_id: str, element_id: str) -> str:
        slide = self.get_slide(slide_id)
        slide["elements"] = [element for element in slide["elements"] if element["id"] != element_id]
        return element_id

    def set_element_style(self, slide_id: str, element_id: str, style: dict[str, Any]) -> str:
        element = self._find_element(slide_id, element_id)
        element["style"] = {**element.get("style", {}), **style}
        return element_id

    def set_background(self, slide_id: str, fill: str | None) -> str:
        self.get_slide(slide_id)["background"] = fill
        return slide_id

    def add_notes(self, slide_id: str, notes_text: str) -> str:
        self.get_slide(slide_id)["speaker_notes"] = notes_text
        return slide_id

    def _find_element(self, slide_id: str, element_id: str) -> dict[str, Any]:
        slide = self.get_slide(slide_id)
        for element in slide["elements"]:
            if element["id"] == element_id:
                return element
        raise KeyError(f"元素不存在：{element_id}")

    def apply_template(self, template_id: str):
        self.template = resolve_ppt_template(template_id)
        self.design_system = design_system_for(self.template)

    # ---------- 与 PPTContent 互转 ----------
    def from_ppt_content(self, content: dict[str, Any]):
        """从既有 PPTContent 恢复为可编辑模型（修订路径）。"""
        self.apply_template(content.get("theme") or self.template["id"])
        self.slides = []
        self._element_seq = 0
        for slide in content.get("slides") or []:
            elements = [dict(element) for element in (slide.get("elements") or [])]
            for element in elements:
                element_id = str(element.get("id") or "")
                if element_id.startswith("E") and element_id[1:].isdigit():
                    self._element_seq = max(self._element_seq, int(element_id[1:]))
            self.slides.append({
                "id": slide.get("id", f"S{len(self.slides) + 1:02d}"),
                "page_type": slide.get("page_type", "concept"),
                "title": slide.get("title", ""),
                "purpose": slide.get("purpose", ""),
                "body": list(slide.get("body") or []),
                "layout": slide.get("layout", "bullet"),
                "visual_suggestion": slide.get("visual_suggestion", ""),
                "speaker_notes": slide.get("speaker_notes", ""),
                "duration_seconds": slide.get("duration_seconds", 0),
                "script_segment_ids": list(slide.get("script_segment_ids") or []),
                "blocks": list(slide.get("blocks") or []),
                "background": slide.get("background"),
                "elements": elements,
                "render_mode": infer_render_mode(slide),
            })
        return self

    def to_ppt_content(self) -> dict[str, Any]:
        """输出 Schema 合法的 PPTContent，同时保持原语义字段逐字不变。"""
        slides = []
        for slide in self.slides:
            body = [str(value) for value in (slide.get("body") or [])]
            blocks = slide.get("blocks") or []
            # Historical Artifacts may intentionally keep a terse ``body`` and
            # richer structured ``blocks``.  A render round-trip must not
            # rewrite body (or its preservation hash).  Only synthesize body
            # for newly-created pages that truly omitted it.
            if blocks and not body:
                body = semantic_body_texts(slide)
            out = {
                "id": slide["id"],
                "page_type": slide["page_type"],
                "title": slide.get("title", ""),
                "purpose": slide.get("purpose", ""),
                "body": body,
                "layout": slide.get("layout", "bullet"),
                "visual_suggestion": slide.get("visual_suggestion", ""),
                "speaker_notes": slide.get("speaker_notes", ""),
                "duration_seconds": int(slide.get("duration_seconds") or 0),
                "script_segment_ids": list(slide.get("script_segment_ids") or []),
            }
            if blocks:
                out["blocks"] = blocks
            if slide.get("elements"):
                out["elements"] = [dict(element) for element in slide["elements"]]
            out["render_mode"] = infer_render_mode(slide)
            slides.append(out)
        return {"theme": self.template["id"], "slides": slides}

    def geometry_report(self) -> list[dict[str, Any]]:
        """输出每元素几何（供几何 QA 检查越界/重叠/文字溢出）。"""
        report = []
        for slide in self.slides:
            for element in self.render_elements(slide):
                report.append({
                    "slide_id": slide["id"], "element_id": element["id"], "kind": element["kind"],
                    "page_type": slide.get("page_type", "concept"),
                    "x": element["x"], "y": element["y"], "w": element["w"], "h": element["h"],
                    "text": element.get("text", ""),
                    "style": element.get("style", {}),
                })
        return report

    def _semantic_elements(self, slide: dict[str, Any]) -> list[dict[str, Any]]:
        """Materialize semantic content for PPTX/QA without mutating the Artifact."""
        template_id = str(self.template.get("id") or "")
        start_x = 2.45 if template_id == "lessonforge_deck_smart_ai" else (2.2 if template_id == "lessonforge_deck_academic" else 0.9)
        media = [item for item in (slide.get("elements") or []) if item.get("kind") in {"image", "chart"}]
        media_x = min((float(item.get("x") or SLIDE_WIDTH) for item in media), default=SLIDE_WIDTH)
        width = max(3.2, min(SLIDE_WIDTH - start_x - 0.65, media_x - start_x - 0.3))
        title = str(slide.get("title") or "")
        body = semantic_body_texts(slide)
        page_type = str(slide.get("page_type") or "concept")
        elements: list[dict[str, Any]] = [{
            "id": "semantic-title", "kind": "textbox", "role": "title", "content_ref": "title",
            "text": title, "x": start_x, "y": 1.65 if page_type == "cover" else 0.7,
            "w": width, "h": 1.35 if page_type == "cover" else 0.75, "z": 1,
            "style": {"size": 38 if page_type == "cover" else 28, "bold": True, "color": "primary"},
        }]
        if body:
            elements.append({
                "id": "semantic-body", "kind": "textbox", "role": "body", "content_ref": "body",
                "text": "\n".join(body), "x": start_x, "y": 3.15 if page_type == "cover" else 1.7,
                "w": width, "h": 1.2 if page_type == "cover" else 4.25, "z": 2,
                "style": {"size": 18 if page_type == "cover" else 17, "color": "text"},
            })
        if page_type == "cover" and slide.get("purpose"):
            elements.append({
                "id": "semantic-purpose", "kind": "textbox", "role": "purpose", "content_ref": "purpose",
                "text": str(slide.get("purpose") or ""), "x": start_x, "y": 4.75,
                "w": width, "h": 0.75, "z": 3,
                "style": {"size": 15, "bold": True, "color": "primary"},
            })
        return elements

    def render_elements(self, slide: dict[str, Any]) -> list[dict[str, Any]]:
        mode = infer_render_mode(slide)
        if mode == "absolute":
            return list(slide.get("elements") or [])
        semantic = self._semantic_elements(slide)
        if mode == "hybrid":
            overlays = [item for item in (slide.get("elements") or []) if item.get("kind") in {"image", "chart"}]
            return [*semantic, *overlays]
        return semantic

    # ---------- 渲染 ----------
    def render(self, output: Path) -> Path:
        """按元素几何 + 模板装饰生成真实 PPTX。"""
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH)
        presentation.slide_height = Inches(SLIDE_HEIGHT)
        palette = self.design_system["palette"]
        fonts = self.design_system["typography"]
        decor = self.design_system["decoration"]
        total = max(1, len(self.slides))
        for index, slide in enumerate(self.slides, 1):
            target = presentation.slides.add_slide(presentation.slide_layouts[6])
            page_type = slide.get("page_type", "concept")
            if page_type == "cover":
                self._decorate(target, palette, fonts, decor.get("cover", []), is_cover=True)
            elif page_type in {"summary", "end"} and page_type == "end":
                self._decorate(target, palette, fonts, decor.get("end", []), is_cover=True)
            else:
                self._decorate(target, palette, fonts, decor.get("shell", []), is_cover=False)
            if slide.get("background"):
                fill = _color(palette, slide["background"]) if slide["background"] in palette else _rgb(slide["background"])
                target.background.fill.solid()
                target.background.fill.fore_color.rgb = fill
            # semantic/hybrid/absolute 统一解析为最终渲染元素。
            for element in sorted(self.render_elements(slide), key=lambda item: item.get("z", 0)):
                self._draw_element(target, element, palette, fonts)
            target.notes_slide.notes_text_frame.text = slide.get("speaker_notes", "")
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output)
        Presentation(output)
        return output

    def _decorate(self, slide, palette, fonts, deco: list[tuple], is_cover: bool):
        for spec in deco:
            kind, x, y, w, h, color_key = spec[0], *spec[1:]
            fill = _color(palette, color_key)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL if kind == "oval" else MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.fill.background()
        # 品牌条
        brand_color = _color(palette, "on_primary" if is_cover else "muted")
        _textbox(slide, "LESSONFORGE 微课", 0.5 if not is_cover else 0.95, 6.92, 4.0, 0.32,
                 font=fonts.get("latin", "Arial"), size=9, color=brand_color, bold=True)

    def _draw_element(self, slide, element: dict[str, Any], palette: dict[str, str], fonts: dict[str, str]):
        kind = element["kind"]
        x, y, w, h = element["x"], element["y"], element["w"], element["h"]
        style = element.get("style") or {}
        font = style.get("font") or fonts.get("body", "Microsoft YaHei")
        if kind == "textbox":
            text = element.get("text", "")
            size = style.get("size", 18)
            color = _color(palette, style["color"]) if style.get("color") in palette else _rgb(style.get("color"), "#1A1A1A")
            _textbox(
                slide, text, x, y, w, h,
                font=font, size=float(size), color=color,
                bold=bool(style.get("bold", False)),
                align=_pp_align(style.get("align")),
                valign=_pp_anchor(style.get("valign")),
            )
        elif kind == "shape":
            fill = _color(palette, element["fill"]) if element.get("fill") in palette else _rgb(element.get("fill"), "#FFFFFF")
            line = _color(palette, element["line"]) if element.get("line") in palette else (_rgb(element["line"]) if element.get("line") else None)
            _shape(slide, element.get("shape_type", "rect"), x, y, w, h, fill, line, radius=bool(element.get("radius", False)))
        elif kind == "image":
            path = Path(element.get("asset_path", ""))
            if path.is_file():
                try:
                    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
                except Exception:  # noqa: BLE001 图片损坏时降级为占位色块
                    _shape(slide, "rect", x, y, w, h, _color(palette, "secondary"), None, radius=True)
            else:
                _shape(slide, "rect", x, y, w, h, _color(palette, "secondary"), None, radius=True)
                _textbox(slide, "图片占位", x + 0.2, y + (h - 0.4) / 2, w - 0.4, 0.4,
                         font=font, size=14, color=_color(palette, "muted"), align=PP_ALIGN.CENTER)
        elif kind == "chart":
            self._draw_chart(slide, element, palette, font)
        elif kind == "note":
            _textbox(slide, element.get("text", ""), x, y, w, h, font=font, size=13,
                     color=_color(palette, "muted"))

    def _draw_chart(self, slide, element: dict[str, Any], palette, font):
        """优先 pptx 原生图表；数据异常时降级为 PIL PNG。"""
        chart = element.get("chart") or {}
        data = chart.get("data") or {}
        categories = data.get("categories") or []
        series = data.get("series") or []
        if not categories or not series:
            _shape(slide, "rect", element["x"], element["y"], element["w"], element["h"],
                   _color(palette, "surface"), _color(palette, "secondary"), radius=True)
            _textbox(slide, "图表占位", element["x"] + 0.3, element["y"] + 0.4, element["w"] - 0.6, 0.5,
                     font=font, size=16, color=_color(palette, "muted"))
            return
        try:
            chart_type = chart.get("type", "bar")
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            chart_data = CategoryChartData()
            chart_data.categories = [str(item) for item in categories]
            for series_item in series:
                chart_data.add_series(series_item.get("name", "系列"), [float(v) for v in series_item.get("values", [])])
            mapping = {"bar": XL_CHART_TYPE.COLUMN_CLUSTERED, "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                       "line": XL_CHART_TYPE.LINE_MARKERS, "pie": XL_CHART_TYPE.PIE}
            chart_kind = mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            frame = slide.shapes.add_chart(
                chart_kind,
                Inches(element["x"]), Inches(element["y"]), Inches(element["w"]), Inches(element["h"]),
                chart_data,
            )
            _style_chart(frame, palette)
        except Exception:  # noqa: BLE001
            try:
                png = render_chart_png(chart.get("type", "bar"), data, palette, (int(element["w"] * 96), int(element["h"] * 96)))
                slide.shapes.add_picture(str(png), Inches(element["x"]), Inches(element["y"]), Inches(element["w"]), Inches(element["h"]))
            except Exception:  # noqa: BLE001
                _shape(slide, "rect", element["x"], element["y"], element["w"], element["h"],
                       _color(palette, "surface"), _color(palette, "secondary"), radius=True)


def _style_chart(frame, palette):
    try:
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart = frame.chart
        chart.has_legend = True
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        except Exception:  # noqa: BLE001
            pass
        try:
            chart.font.size = Pt(11)
        except Exception:  # noqa: BLE001
            pass
    except Exception:  # noqa: BLE001
        pass


def _pp_align(value: str | None):
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(value, PP_ALIGN.LEFT)


def _pp_anchor(value: str | None):
    return {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(value, MSO_ANCHOR.TOP)


def _textbox(slide, text, x, y, w, h, *, font="Microsoft YaHei", size=18, color=RGBColor(26, 26, 26),
             bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _shape(slide, shape_type, x, y, w, h, fill, line=None, radius=False):
    kind = {
        "rect": MSO_SHAPE.RECTANGLE, "oval": MSO_SHAPE.OVAL,
        "rounded": MSO_SHAPE.ROUNDED_RECTANGLE, "line": MSO_SHAPE.RECTANGLE,
    }.get(shape_type, MSO_SHAPE.RECTANGLE)
    if radius:
        kind = MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def render_chart_png(chart_type: str, data: dict[str, Any], palette: dict[str, str],
                     size: tuple[int, int] = (960, 540), output: Path | None = None) -> Path:
    """用 PIL 生成柱/线/饼图 PNG（matplotlib/numpy 缺失时的兜底）。"""
    from app.agent.charting import render_chart_png as _impl
    return _impl(chart_type, data, palette, size, output)
