import hashlib
import re
from pathlib import Path
from uuid import uuid4

from docx import Document
from fastapi import HTTPException, UploadFile
from pypdf import PdfReader
from pptx import Presentation

ALLOWED = {".pdf", ".docx", ".pptx", ".txt", ".md", ".markdown"}


def safe_filename(name: str) -> str:
    base = Path(name).name
    cleaned = re.sub(r"[^\w\u4e00-\u9fff.()（） -]", "_", base).strip(" .")
    if not cleaned or cleaned in {".", ".."}:
        cleaned = "material"
    return cleaned[:180]


async def save_upload(upload: UploadFile, target_dir: Path, max_bytes: int) -> tuple[Path, int, str]:
    suffix = Path(upload.filename or "").suffix.lower()
    if suffix not in ALLOWED:
        raise HTTPException(status_code=415, detail="仅支持 PDF、DOCX、PPTX、TXT 和 Markdown")
    data = await upload.read(max_bytes + 1)
    if len(data) > max_bytes:
        raise HTTPException(status_code=413, detail="文件超过大小限制")
    target_dir.mkdir(parents=True, exist_ok=True)
    path = target_dir / f"{uuid4().hex}{suffix}"
    path.write_bytes(data)
    return path, len(data), hashlib.sha256(data).hexdigest()


def extract_text(path: Path) -> tuple[str, list[dict]]:
    suffix = path.suffix.lower()
    chunks: list[dict] = []
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        for index, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append({"content": text.strip(), "page_number": index, "heading_path": f"第 {index} 页"})
        if not chunks:
            raise ValueError("未提取到文本，该 PDF 可能是扫描件；MVP 暂不支持 OCR")
    elif suffix == ".docx":
        doc = Document(str(path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            paragraphs.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
        chunks = chunk_text("\n".join(paragraphs))
    elif suffix == ".pptx":
        presentation = Presentation(str(path))
        for index, slide in enumerate(presentation.slides, 1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                chunks.append({"content": "\n".join(texts), "page_number": index, "heading_path": f"幻灯片 {index}"})
    else:
        chunks = chunk_text(path.read_text(encoding="utf-8", errors="replace"))
    return "\n\n".join(chunk["content"] for chunk in chunks), chunks


def chunk_text(text: str, size: int = 1800) -> list[dict]:
    cleaned = text.strip()
    return [
        {"content": cleaned[i:i + size], "page_number": None, "heading_path": ""}
        for i in range(0, len(cleaned), size)
    ] if cleaned else []

