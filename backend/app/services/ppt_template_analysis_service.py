"""Deep, deterministic analysis of actual PPTX template files."""
from __future__ import annotations

import hashlib
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.services.ppt_template_service import CATALOG_PATH, ppt_template_catalog_version, resolve_ppt_template


def _font_names(shape) -> list[str]:
    names: list[str] = []
    if not getattr(shape, "has_text_frame", False):
        return names
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                names.append(run.font.name)
    return names


def analyze_template(template_id: str | None) -> tuple[str, dict[str, Any]]:
    metadata = resolve_ppt_template(template_id)
    path = (CATALOG_PATH.parent / metadata["file"]).resolve()
    digest = hashlib.sha256(path.read_bytes()).hexdigest()
    deck = Presentation(str(path))
    fonts: Counter[str] = Counter()
    shape_types: Counter[str] = Counter()
    examples: list[dict[str, Any]] = []
    for index, slide in enumerate(deck.slides):
        if index >= 8:
            break
        items = []
        for shape in slide.shapes:
            kind = str(getattr(shape, "shape_type", "unknown"))
            shape_types[kind] += 1
            fonts.update(_font_names(shape))
            items.append({
                "kind": kind,
                "x": round(shape.left / 914400, 3), "y": round(shape.top / 914400, 3),
                "w": round(shape.width / 914400, 3), "h": round(shape.height / 914400, 3),
                "name": shape.name,
            })
        examples.append({"page": index + 1, "shape_count": len(items), "elements": items[:30]})
    layouts = []
    for layout in deck.slide_layouts:
        placeholders = []
        for shape in layout.placeholders:
            placeholders.append({"name": shape.name, "type": str(shape.placeholder_format.type)})
        layouts.append({"name": layout.name, "placeholders": placeholders})
        fonts.update(name for shape in layout.shapes for name in _font_names(shape))
    profile = {
        "template_id": metadata["id"], "template_hash": digest,
        "catalog_version": ppt_template_catalog_version(), "source_file": path.name,
        "canvas": {"width": round(deck.slide_width / 914400, 3), "height": round(deck.slide_height / 914400, 3)},
        "color_system": metadata.get("palette", {}),
        "palette": metadata.get("palette", {}),
        "typography": {**metadata.get("typography", {}), "observed_fonts": [name for name, _ in fonts.most_common(12)]},
        "masters": len(deck.slide_masters), "layouts": layouts,
        "shape_language": [{"type": name, "count": count} for name, count in shape_types.most_common()],
        "layout_patterns": examples,
        "visual_density": "high" if sum(shape_types.values()) / max(1, len(examples)) > 12 else "medium",
        "design_context_only": True,
    }
    return digest, profile
