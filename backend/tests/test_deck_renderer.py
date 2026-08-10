from pathlib import Path

from app.agents.generators import make_blueprint, make_deck, make_ppt, deck_from_artifact
from app.models.entities import CourseProject
from app.renderers.deck_renderer import deck_template_path, render_deck, role_order, slot_counts


def _course():
    return CourseProject(
        owner_id="u", title="浮力原理", subject="初中物理", grade_level="八年级",
        audience="学生", duration_minutes=15, scenario="课堂讲解",
        language="中文", settings_json={"key_points": "浮力产生的原因与阿基米德原理"},
    )


def test_make_deck_adapts_body_to_template_slot_counts():
    bp = make_blueprint(_course())
    order = role_order()
    assert len(order) == 15
    # 不同模板槽位数不同，make_deck 按模板整形内容
    for template_id in ("lessonforge_deck_academic", "lessonforge_deck_ai_future", "lessonforge_deck_business"):
        deck = make_deck(bp, template_id)
        assert len(deck) == 15
        counts = slot_counts(template_id)
        for index, role in enumerate(order):
            assert len(deck[index]["body"]) == counts.get(role, 0), f"{template_id} {role}"


def test_deck_from_artifact_uses_ai_slides_and_falls_back_to_make_deck():
    bp = make_blueprint(_course())
    template_id = "lessonforge_deck_academic"
    ppt = make_ppt(bp, template_id).model_dump()
    deck = deck_from_artifact(bp, ppt, template_id)
    assert len(deck) == 15
    # AI 生成的 slides 覆盖模板对应页
    assert deck[2]["title"] == ppt["slides"][2]["title"]  # objectives 页
    assert deck[0]["title"] == ppt["slides"][0]["title"]  # 封面
    # 旧 7 页 artifact：前 7 页用 AI 内容，其余用 make_deck 兜底
    legacy = {"theme": template_id, "slides": ppt["slides"][:7]}
    deck_legacy = deck_from_artifact(bp, legacy, template_id)
    assert len(deck_legacy) == 15
    assert deck_legacy[0]["title"] == legacy["slides"][0]["title"]
    fallback = make_deck(bp, template_id)
    assert deck_legacy[7]["title"] == fallback[7]["title"]  # 第 8 页走兜底


def test_render_deck_fills_slots_and_preserves_design():
    from pptx import Presentation

    bp = make_blueprint(_course())
    for template_id in ("lessonforge_deck_academic", "lessonforge_deck_business"):
        deck = make_deck(bp, template_id)
        out = Path(f"/tmp/test_deck_{template_id.split('_')[-1]}.pptx")
        render_deck(deck_template_path(template_id), deck, out, template_id)
        prs = Presentation(str(out))
        assert len(prs.slides) == 15
        slide3_texts = [sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame]
        assert any(text == "学习目标" for text in slide3_texts)
        assert any("OBJ" in text for text in slide3_texts)


def test_export_builds_deck_package_for_deck_theme(tmp_path):
    from app.agents.generators import make_blueprint, make_ppt
    from app.services.export_service import build_course_package
    from pptx import Presentation

    course = _course()
    bp = make_blueprint(course).model_dump()
    ppt = make_ppt(make_blueprint(course)).model_dump()
    ppt["theme"] = "lessonforge_deck_academic"
    build_course_package(
        "c1", "浮力原理", bp, 1,
        {"ppt": {"content_json": ppt, "version": 1}}, tmp_path,
    )
    pptx = tmp_path / "浮力原理_微课资源包" / "02_课件.pptx"
    assert pptx.exists()
    prs = Presentation(str(pptx))
    assert len(prs.slides) == 15
    slide3_texts = [sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame]
    assert any(text == "本课学习目标" for text in slide3_texts)
    assert any("OBJ" in text for text in slide3_texts)


def test_deck_structure_exposes_roles_and_slot_counts():
    from app.renderers.deck_renderer import deck_structure

    structure = deck_structure("lessonforge_deck_academic")
    assert structure["page_count"] == 15
    assert structure["template_id"] == "lessonforge_deck_academic"
    assert len(structure["roles"]) == 15
    roles = {r["role"]: r for r in structure["roles"]}
    assert roles["cover"]["page_type"] == "cover"
    assert roles["cover"]["slot_count"] == slot_counts("lessonforge_deck_academic")["cover"]
    assert roles["objectives"]["page_type"] == "objectives"
    assert roles["end"]["page_type"] == "summary"
    for role, item in roles.items():
        assert item["index"] >= 1 and item["index"] <= 15
        assert item["slot_count"] == slot_counts("lessonforge_deck_academic")[role]


def test_design_system_includes_deck_structure_for_deck_templates():
    from app.renderers.presentation_builder import design_system_for
    from app.services.ppt_template_service import resolve_ppt_template

    deck_template = resolve_ppt_template("lessonforge_deck_academic")
    design = design_system_for(deck_template)
    assert design["composition"] == "deck"
    assert design["deck_structure"]["page_count"] == 15
    # 非 deck 模板不携带 deck_structure（目录当前全为 deck，用合成最小模板测负路径）
    semantic = {"id": "synthetic_semantic", "composition": "semantic"}
    design_semantic = design_system_for(semantic)
    assert design_semantic["composition"] == "semantic"
    assert "deck_structure" not in design_semantic


def test_align_initial_deck_partial_output_maps_by_id_only():
    from types import SimpleNamespace

    from app.agent.pipeline import _align_initial_deck
    from app.renderers.deck_renderer import ROLE_PAGE_TYPE

    bp = make_blueprint(_course())
    runtime = SimpleNamespace(preferred_template="lessonforge_deck_academic", blueprint=bp)
    # 模型只返回 3 页（S05/S10/S15）→ 非 15 页，不做位置回填，缺页走 make_deck 兜底
    partial = [
        {"id": "S05", "title": "模型写的核心页", "body": ["a"], "blocks": [], "speaker_notes": "n"},
        {"id": "S10", "title": "模型写的案例页", "body": ["b"], "blocks": [], "speaker_notes": "n"},
        {"id": "S15", "title": "模型写的末页", "body": ["c"], "blocks": [], "speaker_notes": "n"},
    ]
    aligned = _align_initial_deck(runtime, partial)
    assert len(aligned) == 15
    assert aligned[4]["title"] == "模型写的核心页"   # S05 → 第 5 页
    assert aligned[9]["title"] == "模型写的案例页"    # S10 → 第 10 页
    assert aligned[14]["title"] == "模型写的末页"     # S15 → 第 15 页
    # 缺页用 make_deck 对应角色兜底（封面不被部分页内容错配）
    assert aligned[0]["title"] == make_deck(bp, "lessonforge_deck_academic")[0]["title"]
    assert aligned[0]["page_type"] == "cover"
    # 每页 page_type 与角色对齐
    for index, role in enumerate(role_order()):
        assert aligned[index]["page_type"] == ROLE_PAGE_TYPE[role]


def test_align_initial_deck_order_fallback_only_when_full():
    from types import SimpleNamespace

    from app.agent.pipeline import _align_initial_deck

    bp = make_blueprint(_course())
    runtime = SimpleNamespace(preferred_template="lessonforge_deck_academic", blueprint=bp)
    # 恰好 15 页且无 id/编号 → 位置回填为页序
    full = [{"title": f"页{i + 1}", "body": [], "blocks": [], "speaker_notes": "n"} for i in range(15)]
    aligned = _align_initial_deck(runtime, full)
    assert len(aligned) == 15
    assert aligned[0]["title"] == "页1"
    assert aligned[14]["title"] == "页15"
    assert aligned[0]["page_type"] == "cover"
    assert aligned[14]["page_type"] == "summary"
    # 18 页超量且无 id → 只前 15 页按位置对齐
    too_many = [{"title": f"页{i + 1}", "body": [], "blocks": [], "speaker_notes": "n"} for i in range(18)]
    aligned_many = _align_initial_deck(runtime, too_many)
    assert len(aligned_many) == 15
    assert aligned_many[0]["title"] == "页1"
