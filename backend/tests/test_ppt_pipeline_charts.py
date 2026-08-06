"""PPT 图表工具测试：python-pptx 原生图表 + PIL PNG 兜底。"""
from pathlib import Path

import pytest
from pptx import Presentation

from app.agent.charting import render_chart_png, render_diagram_png
from app.renderers.presentation_builder import PresentationBuilder

PALETTE = {
    "background": "#FFFFFF", "surface": "#F4F7FB", "primary": "#1F4E79",
    "secondary": "#D6E4F0", "text": "#1A1A1A", "muted": "#6B7280", "on_primary": "#FFFFFF",
}


def test_add_chart_creates_native_chart_part(tmp_path):
    builder = PresentationBuilder("lessonforge_deck_academic")
    slide_id = builder.create_slide("comparison", "对比", "comparison")
    builder.add_chart(slide_id, "bar", {
        "categories": ["目标", "现状"],
        "series": [{"name": "对比", "values": [80, 45]}],
    }, 1.0, 1.5, 6.0, 4.0)
    output = Path(tmp_path) / "chart.pptx"
    builder.render(output)
    prs = Presentation(str(output))
    slide = prs.slides[0]
    has_chart = any(shape.has_chart for shape in slide.shapes)
    assert has_chart


def test_add_chart_with_empty_data_degrades_gracefully(tmp_path):
    builder = PresentationBuilder("lessonforge_deck_academic")
    slide_id = builder.create_slide("concept", "标题", "bullet")
    builder.add_chart(slide_id, "pie", {"categories": [], "series": []}, 1.0, 1.5, 4.0, 3.0)
    output = Path(tmp_path) / "empty_chart.pptx"
    builder.render(output)  # 不应抛异常
    assert output.is_file()


@pytest.mark.parametrize("chart_type", ["bar", "line", "pie"])
def test_render_chart_png_produces_image(tmp_path, chart_type):
    output = Path(tmp_path) / f"{chart_type}.png"
    render_chart_png(chart_type, {
        "categories": ["一季度", "二季度", "三季度"],
        "series": [{"name": "成绩", "values": [70, 85, 92]}],
    }, PALETTE, (480, 320), output)
    assert output.is_file()
    from PIL import Image
    with Image.open(output) as image:
        assert image.size == (480, 320)


def test_render_diagram_png_produces_image(tmp_path):
    output = Path(tmp_path) / "flow.png"
    render_diagram_png("flow", {
        "nodes": [{"id": "n1", "label": "问题"}, {"id": "n2", "label": "方法"}, {"id": "n3", "label": "检查"}],
        "edges": [["n1", "n2", "使用"], ["n2", "n3", "验证"]],
    }, PALETTE, (480, 320), output)
    assert output.is_file()
