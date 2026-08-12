"""PresentationBuilder 动态渲染器单测：坐标/元素/round-trip/主题应用。"""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from PIL import Image

from app.renderers.presentation_builder import PresentationBuilder, EMU_PER_INCH
from app.agent.slide_rendering import semantic_content_hash

TEMPLATE = "lessonforge_deck_academic"


@pytest.mark.asyncio
async def test_builder_creates_slides_and_elements():
    builder = PresentationBuilder(TEMPLATE)
    slide_id = builder.create_slide("objectives", "本课学习目标", "bullet", "目标")
    assert slide_id == "S01"
    element_id = builder.add_textbox(slide_id, "要点一", 0.7, 1.8, 5.0, 1.0, style={"size": 18, "color": "text"})
    assert element_id.startswith("E")
    builder.add_shape(slide_id, "rect", 0.2, 0.2, 1.0, 1.0, fill="primary")
    builder.move_element(slide_id, element_id, 1.0, 2.0)
    builder.resize_element(slide_id, element_id, 6.0, 1.2)
    builder.delete_element(slide_id, element_id)
    assert len(builder.get_slide(slide_id)["elements"]) == 1


@pytest.mark.asyncio
async def test_builder_render_produces_valid_pptx_with_exact_coordinates(tmp_path):
    builder = PresentationBuilder(TEMPLATE)
    slide_id = builder.create_slide("objectives", "本课学习目标", "bullet", "目标")
    builder.add_textbox(slide_id, "要点一", 1.0, 2.0, 5.0, 1.0, style={"size": 18, "color": "text"})
    output = Path(tmp_path) / "deck.pptx"
    builder.render(output)
    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    shapes = list(prs.slides[0].shapes)
    text_shapes = [shape for shape in shapes if shape.has_text_frame and shape.text_frame.text == "要点一"]
    assert text_shapes
    box = text_shapes[0]
    assert box.left == Emu(int(1.0 * EMU_PER_INCH))
    assert box.top == Emu(int(2.0 * EMU_PER_INCH))
    assert box.width == Emu(int(5.0 * EMU_PER_INCH))


@pytest.mark.asyncio
async def test_builder_to_ppt_content_round_trip():
    content = {
        "theme": TEMPLATE,
        "slides": [
            {"id": "S01", "page_type": "cover", "title": "浮力原理", "purpose": "p", "body": ["初中物理"],
             "layout": "cover", "visual_suggestion": "封面留白大标题", "speaker_notes": "开场引入课程主题。",
             "duration_seconds": 30, "blocks": []},
            {"id": "S02", "page_type": "objectives", "title": "本课学习目标", "purpose": "p", "body": ["完成本节后能够说明浮力"],
             "layout": "bullet", "visual_suggestion": "编号列表", "speaker_notes": "逐条说明学习目标。",
             "duration_seconds": 60, "blocks": []},
        ],
    }
    builder = PresentationBuilder().from_ppt_content(content)
    out = builder.to_ppt_content()
    assert out["theme"] == TEMPLATE
    assert len(out["slides"]) == 2
    assert out["slides"][0]["title"] == "浮力原理"
    assert out["slides"][0]["duration_seconds"] == 30
    assert out["slides"][1]["duration_seconds"] == 60


def test_builder_round_trip_keeps_untouched_legacy_slide_shape_exact():
    slide = {
        "id": "S01", "page_type": "cover", "title": "浮力原理", "purpose": "p",
        "body": ["初中物理"], "layout": "cover", "visual_suggestion": "封面",
        "speaker_notes": "开场", "duration_seconds": 30,
        "script_segment_ids": [], "blocks": [], "elements": [], "render_mode": None,
    }

    output = PresentationBuilder().from_ppt_content({
        "theme": TEMPLATE, "slides": [slide],
    }).to_ppt_content()["slides"][0]

    assert output == slide


def test_builder_round_trip_keeps_element_ids_unique():
    builder = PresentationBuilder().from_ppt_content({
        "theme": TEMPLATE,
        "slides": [{
            "id": "S01", "page_type": "cover", "title": "浮力",
            "elements": [{
                "id": "E09", "kind": "image", "x": 7, "y": 1, "w": 5, "h": 4,
                "z": 9, "style": {}, "asset_path": "/tmp/existing.png", "role": "visual",
            }],
        }],
    })
    new_id = builder.add_textbox("S01", "新模板标题", 1, 1, 5, 1)
    assert new_id == "E10"
    assert len({item["id"] for item in builder.get_slide("S01")["elements"]}) == 2


@pytest.mark.asyncio
async def test_apply_template_switches_design_system():
    builder = PresentationBuilder(TEMPLATE)
    before = builder.design_system["id"]
    builder.apply_template("lessonforge_deck_smart_ai")
    assert builder.design_system["id"] == "lessonforge_deck_smart_ai"
    assert builder.design_system["palette"]["primary"] == "#5B3DF5"
    assert before == TEMPLATE


def test_geometry_report_lists_elements():
    builder = PresentationBuilder(TEMPLATE)
    slide_id = builder.create_slide("concept", "标题", "bullet")
    builder.add_textbox(slide_id, "正文", 0.7, 1.8, 5.0, 1.0)
    report = builder.geometry_report()
    assert len(report) == 1
    assert report[0]["slide_id"] == slide_id
    assert report[0]["x"] == 0.7


def test_historical_media_only_page_is_inferred_as_hybrid_and_keeps_semantic_text(tmp_path):
    image_path = tmp_path / "buoyancy.png"
    Image.new("RGB", (640, 480), "navy").save(image_path)
    content = {
        "theme": TEMPLATE,
        "slides": [{
            "id": "slide_03_km", "page_type": "concept",
            "title": "从液体压强推导至阿基米德原理",
            "purpose": "推导浮力来源", "body": ["左右表面压力抵消", "上下表面存在压力差"],
            "layout": "bullet", "visual_suggestion": "潜水艇浮力图",
            "speaker_notes": "保留教师讲解备注", "duration_seconds": 60,
            "blocks": [{"kind": "bullets", "items": [{"text": "压力差产生向上的浮力"}]}],
            # V34 历史结构：有语义内容，但 elements 只有一张图片且没有 render_mode。
            "elements": [{
                "id": "E34", "kind": "image", "role": "visual", "visual_slot": "primary_visual",
                "x": 7.7, "y": 1.5, "w": 4.8, "h": 3.6, "asset_path": str(image_path),
                "asset_id": "asset-v34", "style": {},
            }],
        }],
    }
    builder = PresentationBuilder().from_ppt_content(content)
    slide = builder.get_slide("slide_03_km")
    assert slide["render_mode"] == "hybrid"
    rendered = builder.render_elements(slide)
    assert any(item.get("content_ref") == "title" for item in rendered)
    assert any(item.get("kind") == "image" for item in rendered)

    output = tmp_path / "v34-hybrid.pptx"
    builder.render(output)
    presentation = Presentation(str(output))
    texts = "\n".join(shape.text for shape in presentation.slides[0].shapes if getattr(shape, "has_text_frame", False))
    assert "从液体压强推导至阿基米德原理" in texts
    assert "压力差产生向上的浮力" in texts
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in presentation.slides[0].shapes)


def test_builder_roundtrip_preserves_hash_when_body_and_blocks_are_distinct():
    slide = {
        "id": "slide_04", "page_type": "concept", "title": "浮力推导", "purpose": "解释来源",
        "body": ["原始摘要不得被重算"],
        "blocks": [{"kind": "steps", "steps": [{"title": "比较压力", "detail": "上下表面压力差形成浮力"}]}],
        "layout": "steps", "visual_suggestion": "受力图", "speaker_notes": "保留原备注",
        "duration_seconds": 60, "script_segment_ids": ["seg-04"],
    }
    before = semantic_content_hash(slide)
    output = PresentationBuilder().from_ppt_content({"theme": TEMPLATE, "slides": [slide]}).to_ppt_content()["slides"][0]
    assert output["body"] == ["原始摘要不得被重算"]
    assert output["script_segment_ids"] == ["seg-04"]
    assert semantic_content_hash(output) == before


def test_hybrid_pptx_renders_step_details_quote_citation_and_notes(tmp_path):
    image_path = tmp_path / "visual.png"
    Image.new("RGB", (640, 480), "teal").save(image_path)
    content = {"theme": TEMPLATE, "slides": [{
        "id": "slide_04", "page_type": "concept", "title": "阿基米德原理", "purpose": "",
        "body": ["结构化正文摘要"], "layout": "steps", "visual_suggestion": "",
        "speaker_notes": "逐步说明压力差与排液重力的关系。", "duration_seconds": 60,
        "blocks": [
            {"kind": "steps", "steps": [{"title": "比较压力", "detail": "p=ρgh，上下表面压力不同"}]},
            {"kind": "quote", "text": "F浮=G排", "citation": "阿基米德原理"},
        ],
        "elements": [{
            "id": "E01", "kind": "image", "role": "visual", "visual_slot": "primary_visual",
            "x": 7.7, "y": 1.5, "w": 4.8, "h": 3.6, "asset_path": str(image_path), "asset_id": "asset-1",
        }],
    }]}
    output = tmp_path / "structured-hybrid.pptx"
    PresentationBuilder().from_ppt_content(content).render(output)
    presentation = Presentation(str(output))
    texts = "\n".join(shape.text for shape in presentation.slides[0].shapes if getattr(shape, "has_text_frame", False))
    assert "p=ρgh，上下表面压力不同" in texts
    assert "阿基米德原理" in texts
    assert "F浮=G排" in texts
    assert "逐步说明压力差与排液重力的关系。" in presentation.slides[0].notes_slide.notes_text_frame.text
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in presentation.slides[0].shapes)
