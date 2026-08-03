import zipfile
from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from pydantic import ValidationError
from pptx import Presentation

from app.agents.generators import make_blueprint, make_exercises, make_lesson_plan, make_ppt, make_task_sheet, make_verbatim, make_video_script, to_markdown
from app.models.entities import CourseProject
from app.renderers.docx_renderer import render_markdown_docx, render_task_sheet_docx
from app.schemas.artifact import TaskSheetContent
from app.renderers.pptx_renderer import render_pptx
from app.services.export_service import build_course_package, safe_package_name
from app.services.material_service import safe_filename
from app.services.ppt_template_service import list_ppt_templates
from app.services.quality_service import estimate_chinese_minutes, validate_resources, validate_task_sheet


def sample_course():
    return CourseProject(owner_id="u", title="牛顿第二定律", subject="高中物理", grade_level="高一", audience="已学习运动学基础的学生", duration_minutes=15, scenario="课堂讲解", language="中文", settings_json={"course_task": "解释力、质量与加速度的关系"})


def test_schemas_rules_and_timing():
    bp = make_blueprint(sample_course())
    ppt, exercise = make_ppt(bp), make_exercises(bp)
    lesson = make_lesson_plan(bp)
    script = make_video_script(bp, lesson, ppt)
    verbatim = make_verbatim(bp, ppt, script)
    data = {"lesson_plan": lesson.model_dump(), "ppt": ppt.model_dump(), "exercise": exercise.model_dump(), "video_script": script.model_dump(), "verbatim": verbatim.model_dump()}
    assert validate_resources(bp, data) == []
    assert sum(s.end_minute - s.start_minute for s in bp.timeline) == 15
    assert estimate_chinese_minutes("教" * 440) == 2


def test_safe_filenames():
    assert safe_filename("../../教案?.pdf").endswith(".pdf")
    assert "/" not in safe_package_name("物理/力学:第一课")


def test_docx_pptx_and_zip(tmp_path: Path):
    bp = make_blueprint(sample_course())
    values = {"lesson_plan": make_lesson_plan(bp), "ppt": make_ppt(bp), "task_sheet": make_task_sheet(bp), "exercise": make_exercises(bp)}
    values["video_script"] = make_video_script(bp, values["lesson_plan"], values["ppt"])
    values["verbatim"] = make_verbatim(bp, values["ppt"], values["video_script"])
    doc = render_markdown_docx("教学设计", to_markdown("lesson_plan", values["lesson_plan"]), tmp_path / "test.docx")
    ppt = render_pptx("测试课程", values["ppt"].model_dump(), tmp_path / "test.pptx")
    assert len(Document(doc).paragraphs) > 3
    assert len(Presentation(ppt).slides) == len(values["ppt"].slides)
    artifacts = {key: {"version": 1, "content_json": value.model_dump(), "content_markdown": to_markdown(key, value)} for key, value in values.items()}
    package, manifest = build_course_package("course", "测试课程", bp.model_dump(), 1, artifacts, tmp_path / "out")
    assert package.exists() and len(manifest["artifacts"]) >= 9
    ppt_manifest = next(item for item in manifest["artifacts"] if item["type"] == "ppt")
    assert ppt_manifest["template_id"] == "lessonforge_swiss_blue"
    assert ppt_manifest["template_catalog_version"]
    with zipfile.ZipFile(package) as archive:
        assert archive.testzip() is None
        names = archive.namelist()
        task_docx_name = next(name for name in names if name.endswith("03_学习任务单.docx"))
        task_markdown_name = next(name for name in names if name.endswith("03_学习任务单.md"))
        video_docx_name = next(name for name in names if name.endswith("06_微课视频脚本.docx"))
        video_markdown_name = next(name for name in names if name.endswith("06_微课视频脚本.md"))
        task_document = Document(BytesIO(archive.read(task_docx_name)))
        markdown = archive.read(task_markdown_name).decode("utf-8")
        assert len(task_document.tables) >= 6
        assert "学习观察记录" in markdown
        assert "参考答案" not in markdown and "教师提示" not in markdown
        assert len(Document(BytesIO(archive.read(video_docx_name))).tables) >= 2
        assert "旁白" in archive.read(video_markdown_name).decode("utf-8")


def test_all_ppt_templates_render_and_reopen(tmp_path: Path):
    bp = make_blueprint(sample_course())
    base = make_ppt(bp).model_dump()
    for template in list_ppt_templates():
        content = {**base, "theme": template["id"]}
        path = render_pptx("模板测试", content, tmp_path / f"{template['id']}.pptx")
        presentation = Presentation(path)
        assert len(presentation.slides) == len(base["slides"])
        assert presentation.slide_width == 12191695
        assert presentation.slide_height == 6858000
        assert presentation.slides[0].notes_slide.notes_text_frame.text
        assert all(len(slide.shapes) >= 5 for slide in presentation.slides)


def test_task_sheet_v2_schema_and_quality_paths(tmp_path: Path):
    bp = make_blueprint(sample_course())
    sheet = make_task_sheet(bp)
    payload = sheet.model_dump()
    assert TaskSheetContent.model_validate(payload).schema_version == "2.0"
    assert validate_task_sheet(bp, payload, make_lesson_plan(bp).model_dump()) == []

    missing_output = sheet.model_dump()
    missing_output["tasks"][0].pop("student_output")
    with pytest.raises(ValidationError):
        TaskSheetContent.model_validate(missing_output)

    invalid = sheet.model_copy(deep=True).model_dump()
    invalid["learning_objectives"][0]["id"] = "OBJ-404"
    invalid["tasks"][0]["objective_ids"] = ["OBJ-404"]
    invalid["tasks"][0]["knowledge_point_ids"] = ["KP-404"]
    invalid["tasks"][0]["estimated_minutes"] = 99
    invalid["learning_questions"][1]["id"] = invalid["learning_questions"][0]["id"]
    invalid["record_table"] = None
    invalid["tasks"] = [{**task, "record_table": None} for task in invalid["tasks"]]
    issues = validate_task_sheet(bp, invalid, make_lesson_plan(bp).model_dump())
    assert issues
    assert all(item["target_agent"] == "task_sheet_agent" for item in issues)
    assert all(item["location"].startswith("$.") for item in issues)
    descriptions = "\n".join(item["description"] for item in issues)
    assert "OBJ-404" in descriptions
    assert "KP-404" in descriptions
    assert "没有可填写" in descriptions
    assert "超过" in descriptions

    path = render_task_sheet_docx("学习任务单", payload, tmp_path / "task-sheet.docx", "V2")
    document = Document(path)
    assert len(document.tables) >= 6
    assert any("学生版" in paragraph.text for section in document.sections for paragraph in section.footer.paragraphs)


def test_legacy_task_sheet_export_uses_markdown_fallback(tmp_path: Path):
    bp = make_blueprint(sample_course())
    legacy = {
        "version": 7,
        "content_json": {
            "learning_objectives": ["OBJ-01：解释核心概念"],
            "tasks": [{"id": "T-01", "action": "观察", "object": "情境", "output": "记录", "completion_criterion": "填写完整"}],
        },
        "content_markdown": "# 旧版学习任务单\n\n- 观察情境并完成记录",
    }
    package, _ = build_course_package(
        "legacy-course", "旧版课程", bp.model_dump(), 1,
        {"task_sheet": legacy}, tmp_path / "legacy-out",
    )
    with zipfile.ZipFile(package) as archive:
        names = archive.namelist()
        docx_name = next(name for name in names if name.endswith("03_学习任务单.docx"))
        markdown_name = next(name for name in names if name.endswith("03_学习任务单.md"))
        assert Document(BytesIO(archive.read(docx_name))).paragraphs
        assert "旧版学习任务单" in archive.read(markdown_name).decode("utf-8")


def test_markdown_renderers_cover_all_lesson_plan_sections():
    lesson = make_lesson_plan(make_blueprint(sample_course()))
    lesson.reflection_placeholder = "REFLECTION_MARKER"
    lesson.board_design = "BOARD_MARKER"
    lesson.homework = "HOMEWORK_MARKER"

    markdown = to_markdown("lesson_plan", lesson)

    for heading in (
        "内容分析", "学情分析", "教学目标", "教学重点", "教学难点",
        "教学方法与策略", "教学资源", "教学过程", "板书设计", "作业布置", "教学反思",
    ):
        assert f"## {heading}" in markdown
    assert "REFLECTION_MARKER" in markdown
    assert "BOARD_MARKER" in markdown
    assert "HOMEWORK_MARKER" in markdown
    assert all(stage.assessment in markdown for stage in lesson.stages)
