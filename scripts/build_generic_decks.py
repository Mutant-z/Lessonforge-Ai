#!/usr/bin/env python3
"""重建 templates/PPT_template 为 6 套独立设计系统的微课模板，并同步槽位几何。

目标（参考 KIMI 的 PPT 模板理念）：每套模板是独立的视觉设计系统——封面样式、
内容页布局、装饰、排版、槽位几何都不同，而不是只换配色。
6 套设计语言：
  deck_academic         学术科研 · 左栏编号学术
  deck_ai_future        AI 未来   · 深色科技卡片
  deck_business         商务培训 · 极简数据两栏
  deck_cartoon          卡通启蒙 · 圆角徽章卡片
  deck_chinese_culture  中国文化 · 纸墨竖排
  deck_smart_ai         智慧课堂 · 紫色渐变侧栏分栏

每个模板的封面/内容页槽位 (x,y) 即文本框左上角，也是 deck_slots.json 的锚点，
由 deck_renderer 按模板读槽位并填字。

用法：在仓库根目录执行  .venv/bin/python scripts/build_generic_decks.py
"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
CATALOG_PATH = ROOT / "templates" / "pptx" / "catalog.json"
DECK_DIR = ROOT / "templates" / "PPT_template"
SLOTS_PATH = ROOT / "templates" / "ppt_decks" / "deck_slots.json"

EMU_PER_INCH = 914400
CANVAS_W, CANVAS_H = 13.333, 7.5

ROLE_ORDER = [
    "cover", "intro", "objectives", "knowledge_map", "knowledge_intro",
    "core_1", "core_2", "core_3", "core_4", "case_study",
    "discussion", "summary", "assessment", "assignment", "end",
]

# 槽位规格：(x, y, w, h, 字号, 颜色键, 是否加粗, 对齐, 是否斜体)
def _s(x, y, w, h, size, color, bold=False, align="left", italic=False):
    return (x, y, w, h, size, color, bold, align, italic)


def _stack(x, y, w, h, n, gap, size, color, bold=False, align="left"):
    return [_s(x, y + i * (h + gap), w, h, size, color, bold, align) for i in range(n)]


def _grid(x, y, w, h, cols, rows, gap_x, gap_y, size, color, bold=False, align="left"):
    """行优先网格：slot0 左上 → slot1 右上 → …"""
    slots = []
    for row in range(rows):
        for col in range(cols):
            slots.append(_s(x + col * (w + gap_x), y + row * (h + gap_y), w, h, size, color, bold, align))
    return slots


def _twocol(x0, x1, y, w, h, per_col, gap_y, size, color, bold=False, align="left"):
    """左列 per_col 条，右列 per_col 条（行优先）。"""
    slots = []
    for i in range(per_col):
        slots.append(_s(x0, y + i * (h + gap_y), w, h, size, color, bold, align))
        slots.append(_s(x1, y + i * (h + gap_y), w, h, size, color, bold, align))
    return slots


# ============ 6 套模板设计系统 ============
# 每套：cover（整页背景页）、end（整页背景页）、shell（内容页公共）、roles（各角色内容槽位）
# content_slots 的 (x,y) 即槽位锚点，也是 deck_slots.json 的锚点。
# card：若存在，内容槽位会套上圆角卡片（背景 surface + 边框 secondary）。
TEMPLATE_DESIGNS: dict[str, dict] = {
    # ---------- 1. 学术科研：左栏编号学术，单列正式列表 ----------
    "lessonforge_deck_academic": {
        "cover": {
            "bg": "background",
            "deco": [
                ("rect", 0, 0, 2.1, 7.5, "primary"),
                ("rect", 2.1, 0, 0.06, 7.5, "secondary"),
            ],
            "brand": (0.55, 6.9, 1.5, 0.4, "on_primary"),
            "title": _s(2.7, 2.35, 9.6, 1.7, 40, "primary", True),
            "subtitle": _s(2.7, 4.15, 9.6, 0.7, 18, "muted"),
        },
        "end": {
            "bg": "primary",
            "deco": [("oval", -1.5, 4.5, 4.8, 4.8, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.2, 2.4, 11.0, 1.3, 34, "on_primary", True),
            "subtitle": _s(1.2, 3.95, 11.0, 0.8, 20, "on_primary"),
        },
        "shell": {
            "deco": [("rect", 0, 0, 0.16, 7.5, "primary")],
            "brand": (0.5, 7.02, 4, 0.3, "muted"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(0.65, 0.5, 11.8, 0.85, 28, "primary", True),
            "accent": ("rect", 0.68, 1.42, 1.35, 0.09, "secondary"),
        },
        "roles": {
            "intro": {"content": [_s(0.65, 1.85, 12.0, 1.0, 18, "primary", True)] + _stack(0.65, 3.05, 12.0, 0.58, 5, 0.68, 16, "text")},
            "objectives": {"content": [_s(0.65, 1.9, 12.0, 0.65, 18, "primary", True)] + _stack(0.65, 2.7, 12.0, 0.72, 4, 0.85, 16, "text")},
            "knowledge_map": {"content": [_s(0.65, 1.9, 12.0, 0.65, 18, "primary", True)] + _twocol(0.65, 6.85, 2.7, 5.9, 0.72, 2, 0.9, 16, "text")},
            "knowledge_intro": {"content": [_s(0.65, 1.85, 12.0, 0.95, 18, "primary", True)] + _twocol(0.65, 6.85, 2.95, 5.9, 0.62, 3, 0.75, 16, "text")},
            "core_1": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True), _s(0.65, 2.45, 12.0, 0.75, 18, "text", True)] + _stack(0.65, 3.35, 12.0, 0.62, 4, 0.75, 16, "text")},
            "core_2": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True), _s(0.65, 2.45, 12.0, 0.75, 18, "text", True)] + _stack(0.65, 3.35, 12.0, 0.62, 4, 0.75, 16, "text")},
            "core_3": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True), _s(0.65, 2.45, 12.0, 0.75, 18, "text", True)] + _stack(0.65, 3.35, 12.0, 0.62, 4, 0.75, 16, "text")},
            "core_4": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True), _s(0.65, 2.45, 12.0, 0.75, 18, "text", True)] + _stack(0.65, 3.35, 12.0, 0.62, 4, 0.75, 16, "text")},
            "case_study": {"content": _grid(0.65, 1.85, 2.85, 1.1, 4, 1, 3.1, 0, 16, "text", True) + [_s(0.65, 3.35, 3.0, 0.5, 14, "primary", True), _s(0.65, 3.95, 12.0, 1.2, 16, "text")]},
            "discussion": {"content": [_s(0.65, 1.95, 12.0, 0.85, 18, "primary", True)] + _stack(0.65, 3.1, 12.0, 0.8, 3, 1.0, 16, "text")},
            "summary": {"content": [_s(0.65, 1.8, 12.0, 0.6, 17, "primary", True)] + _twocol(0.65, 6.85, 2.55, 5.9, 0.6, 4, 0.75, 16, "text") + [_s(0.65, 5.85, 12.0, 0.7, 15, "muted", False, "left", True)]},
            "assessment": {"content": [_s(1.0, 1.85, 11.3, 0.5, 15, "text", True), _s(1.0, 2.4, 11.3, 0.45, 14, "muted"), _s(1.0, 3.3, 11.3, 0.5, 15, "text", True), _s(1.0, 3.85, 11.3, 0.45, 14, "muted"), _s(1.0, 4.75, 11.3, 0.5, 15, "text", True), _s(1.0, 5.3, 11.3, 0.45, 14, "muted"), _s(1.0, 6.2, 11.3, 0.4, 12, "muted", False, "left", True)]},
            "assignment": {"content": [_s(1.0, 1.9, 6.0, 0.5, 14, "primary", True), _s(1.0, 2.5, 11.3, 1.0, 17, "text", True)] + _twocol(1.0, 6.8, 3.95, 5.5, 0.72, 2, 0.9, 16, "text")},
        },
    },

    # ---------- 2. AI 未来：深色科技卡片，2×2 圆角卡片 ----------
    "lessonforge_deck_ai_future": {
        "cover": {
            "bg": "primary",
            "deco": [
                ("rect", 0, 0, 13.333, 0.1, "secondary"),
                ("rect", 1.0, 1.9, 11.3, 0.02, "secondary"),
                ("rect", 0.6, 6.4, 6.0, 0.02, "secondary"),
            ],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.0, 2.5, 11.3, 1.7, 42, "on_primary", True),
            "subtitle": _s(1.0, 4.4, 11.3, 0.7, 20, "secondary"),
        },
        "end": {
            "bg": "primary",
            "deco": [("rect", 0, 0, 13.333, 0.1, "secondary"), ("oval", 11.4, -1.4, 3.6, 3.6, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.2, 2.4, 11.0, 1.3, 34, "on_primary", True),
            "subtitle": _s(1.2, 3.95, 11.0, 0.8, 20, "secondary"),
        },
        "shell": {
            "deco": [("rect", 0, 0, 13.333, 0.12, "primary"), ("oval", 12.6, 0.35, 0.5, 0.5, "secondary")],
            "brand": (0.5, 7.02, 4, 0.3, "muted"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(0.7, 0.6, 11.6, 0.8, 28, "primary", True),
            "accent": ("rect", 0.73, 1.48, 1.35, 0.08, "secondary"),
            "card": {"pad": 0.28, "fill": "surface", "border": "secondary", "radius": 0.28},
        },
        "roles": {
            "intro": {"content": [_s(0.7, 1.85, 12.0, 1.0, 18, "primary", True)] + _grid(0.7, 3.1, 5.85, 1.6, 2, 2, 0.35, 0.35, 16, "text")},
            "objectives": {"content": _grid(0.7, 1.85, 5.85, 1.95, 2, 2, 0.35, 0.4, 16, "text")},
            "knowledge_map": {"content": _grid(0.7, 1.85, 5.85, 1.95, 2, 2, 0.35, 0.4, 16, "text")},
            "knowledge_intro": {"content": [_s(0.7, 1.85, 12.0, 0.9, 18, "primary", True)] + _grid(0.7, 3.0, 5.85, 1.6, 2, 2, 0.35, 0.35, 16, "text")},
            "core_1": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _grid(0.7, 2.5, 5.85, 1.8, 2, 2, 0.35, 0.35, 16, "text")},
            "core_2": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _grid(0.7, 2.5, 5.85, 1.8, 2, 2, 0.35, 0.35, 16, "text")},
            "core_3": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _grid(0.7, 2.5, 5.85, 1.8, 2, 2, 0.35, 0.35, 16, "text")},
            "core_4": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _grid(0.7, 2.5, 5.85, 1.8, 2, 2, 0.35, 0.35, 16, "text")},
            "case_study": {"content": _grid(0.7, 1.85, 5.85, 1.7, 2, 2, 0.35, 0.35, 16, "text") + [_s(0.7, 5.35, 12.0, 0.9, 16, "text")]},
            "discussion": {"content": [_s(0.7, 1.85, 12.0, 0.85, 18, "primary", True)] + _grid(0.7, 3.0, 5.85, 1.6, 2, 2, 0.35, 0.35, 16, "text")},
            "summary": {"content": [_s(0.7, 1.8, 12.0, 0.6, 17, "primary", True)] + _grid(0.7, 2.6, 5.85, 1.75, 2, 2, 0.35, 0.4, 16, "text")},
            "assessment": {"content": [_s(1.0, 1.85, 11.3, 0.5, 15, "text", True), _s(1.0, 2.4, 11.3, 0.45, 14, "muted"), _s(1.0, 3.3, 11.3, 0.5, 15, "text", True), _s(1.0, 3.85, 11.3, 0.45, 14, "muted"), _s(1.0, 4.75, 11.3, 0.5, 15, "text", True), _s(1.0, 5.3, 11.3, 0.45, 14, "muted"), _s(1.0, 6.2, 11.3, 0.4, 12, "muted", False, "left", True)]},
            "assignment": {"content": [_s(1.0, 1.9, 6.0, 0.5, 14, "primary", True), _s(1.0, 2.5, 11.3, 1.0, 17, "text", True)] + _grid(1.0, 3.8, 5.5, 1.3, 2, 2, 0.3, 0.3, 16, "text")},
        },
    },

    # ---------- 3. 商务培训：极简数据两栏，指标徽标 ----------
    "lessonforge_deck_business": {
        "cover": {
            "bg": "background",
            "deco": [("rect", 0, 0, 13.333, 0.35, "primary"), ("rect", 1.0, 3.95, 2.2, 0.09, "secondary")],
            "brand": (1.0, 6.9, 4, 0.4, "muted"),
            "title": _s(1.0, 2.5, 11.3, 1.4, 44, "primary", True),
            "subtitle": _s(1.0, 4.2, 11.3, 0.8, 22, "muted"),
        },
        "end": {
            "bg": "primary",
            "deco": [("rect", 0, 0, 13.333, 0.35, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.2, 2.4, 11.0, 1.3, 34, "on_primary", True),
            "subtitle": _s(1.2, 3.95, 11.0, 0.8, 20, "secondary"),
        },
        "shell": {
            "deco": [("rect", 0, 0, 13.333, 0.1, "primary"), ("rect", 0.65, 6.85, 12.0, 0.015, "secondary")],
            "brand": (0.5, 7.02, 4, 0.3, "muted"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(0.65, 0.55, 11.8, 0.8, 28, "primary", True),
            "accent": ("rect", 0.68, 1.42, 1.35, 0.09, "primary"),
        },
        "roles": {
            "intro": {"content": [_s(0.65, 1.85, 12.0, 1.0, 18, "primary", True)] + _twocol(0.65, 6.85, 3.1, 5.9, 0.85, 3, 0.5, 16, "text")},
            "objectives": {"content": _twocol(0.65, 6.85, 1.95, 5.9, 1.15, 3, 0.55, 16, "text")},
            "knowledge_map": {"content": _twocol(0.65, 6.85, 1.95, 5.9, 1.15, 3, 0.55, 16, "text")},
            "knowledge_intro": {"content": [_s(0.65, 1.85, 12.0, 0.95, 18, "primary", True)] + _twocol(0.65, 6.85, 3.0, 5.9, 0.8, 3, 0.45, 16, "text")},
            "core_1": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.65, 6.85, 2.5, 5.9, 1.05, 3, 0.5, 16, "text")},
            "core_2": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.65, 6.85, 2.5, 5.9, 1.05, 3, 0.5, 16, "text")},
            "core_3": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.65, 6.85, 2.5, 5.9, 1.05, 3, 0.5, 16, "text")},
            "core_4": {"content": [_s(0.65, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.65, 6.85, 2.5, 5.9, 1.05, 3, 0.5, 16, "text")},
            "case_study": {"content": _grid(0.65, 1.85, 2.85, 1.1, 4, 1, 3.1, 0, 16, "text", True) + [_s(0.65, 3.35, 3.0, 0.5, 14, "primary", True), _s(0.65, 3.95, 12.0, 1.2, 16, "text")]},
            "discussion": {"content": [_s(0.65, 1.95, 12.0, 0.85, 18, "primary", True)] + _twocol(0.65, 6.85, 3.1, 5.9, 0.9, 2, 0.6, 16, "text")},
            "summary": {"content": [_s(0.65, 1.8, 12.0, 0.6, 17, "primary", True)] + _twocol(0.65, 6.85, 2.6, 5.9, 0.95, 3, 0.5, 16, "text") + [_s(0.65, 5.9, 12.0, 0.7, 15, "muted", False, "left", True)]},
            "assessment": {"content": [_s(1.0, 1.85, 11.3, 0.5, 15, "text", True), _s(1.0, 2.4, 11.3, 0.45, 14, "muted"), _s(1.0, 3.3, 11.3, 0.5, 15, "text", True), _s(1.0, 3.85, 11.3, 0.45, 14, "muted"), _s(1.0, 4.75, 11.3, 0.5, 15, "text", True), _s(1.0, 5.3, 11.3, 0.45, 14, "muted"), _s(1.0, 6.2, 11.3, 0.4, 12, "muted", False, "left", True)]},
            "assignment": {"content": [_s(1.0, 1.9, 6.0, 0.5, 14, "primary", True), _s(1.0, 2.5, 11.3, 1.0, 17, "text", True)] + _twocol(1.0, 6.8, 3.95, 5.5, 0.72, 2, 0.9, 16, "text")},
        },
    },

    # ---------- 4. 卡通启蒙：圆角徽章卡片，奶油底 ----------
    "lessonforge_deck_cartoon": {
        "cover": {
            "bg": "background",
            "deco": [
                ("oval", 11.0, 0.4, 1.6, 1.6, "secondary"),
                ("oval", 0.7, 0.5, 1.1, 1.1, "secondary"),
                ("rect", 1.0, 3.6, 3.0, 0.12, "secondary"),
            ],
            "brand": (1.0, 6.9, 4, 0.4, "muted"),
            "title": _s(1.0, 2.3, 11.3, 1.3, 42, "primary", True),
            "subtitle": _s(1.0, 3.95, 11.3, 0.8, 22, "muted"),
        },
        "end": {
            "bg": "primary",
            "deco": [("oval", 11.2, -1.2, 3.2, 3.2, "secondary"), ("oval", -0.8, 5.6, 2.6, 2.6, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.2, 2.4, 11.0, 1.3, 34, "on_primary", True),
            "subtitle": _s(1.2, 3.95, 11.0, 0.8, 20, "secondary"),
        },
        "shell": {
            "deco": [("oval", 12.3, 0.4, 0.7, 0.7, "secondary")],
            "brand": (0.5, 7.02, 4, 0.3, "muted"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(0.7, 0.55, 11.6, 0.85, 30, "primary", True),
            "accent": ("rect", 0.73, 1.48, 1.35, 0.1, "secondary"),
            "card": {"pad": 0.22, "fill": "surface", "border": "secondary", "radius": 0.35},
        },
        "roles": {
            "intro": {"content": [_s(0.7, 1.85, 12.0, 1.0, 18, "primary", True)] + _twocol(0.7, 6.7, 3.1, 5.7, 0.95, 3, 0.45, 17, "text")},
            "objectives": {"content": _twocol(0.7, 6.7, 1.95, 5.7, 1.2, 3, 0.5, 17, "text")},
            "knowledge_map": {"content": _twocol(0.7, 6.7, 1.95, 5.7, 1.2, 3, 0.5, 17, "text")},
            "knowledge_intro": {"content": [_s(0.7, 1.85, 12.0, 0.9, 18, "primary", True)] + _twocol(0.7, 6.7, 3.0, 5.7, 0.9, 3, 0.4, 17, "text")},
            "core_1": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.7, 6.7, 2.5, 5.7, 1.1, 3, 0.45, 17, "text")},
            "core_2": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.7, 6.7, 2.5, 5.7, 1.1, 3, 0.45, 17, "text")},
            "core_3": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.7, 6.7, 2.5, 5.7, 1.1, 3, 0.45, 17, "text")},
            "core_4": {"content": [_s(0.7, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(0.7, 6.7, 2.5, 5.7, 1.1, 3, 0.45, 17, "text")},
            "case_study": {"content": _grid(0.7, 1.85, 2.85, 1.15, 4, 1, 3.1, 0, 17, "text", True) + [_s(0.7, 3.4, 3.0, 0.5, 14, "primary", True), _s(0.7, 4.0, 12.0, 1.2, 17, "text")]},
            "discussion": {"content": [_s(0.7, 1.95, 12.0, 0.85, 18, "primary", True)] + _twocol(0.7, 6.7, 3.1, 5.7, 1.0, 2, 0.6, 17, "text")},
            "summary": {"content": [_s(0.7, 1.8, 12.0, 0.6, 17, "primary", True)] + _twocol(0.7, 6.7, 2.6, 5.7, 1.0, 3, 0.5, 17, "text") + [_s(0.7, 6.0, 12.0, 0.6, 15, "muted", False, "left", True)]},
            "assessment": {"content": [_s(1.0, 1.85, 11.3, 0.55, 16, "text", True), _s(1.0, 2.45, 11.3, 0.5, 15, "muted"), _s(1.0, 3.4, 11.3, 0.55, 16, "text", True), _s(1.0, 4.0, 11.3, 0.5, 15, "muted"), _s(1.0, 4.95, 11.3, 0.55, 16, "text", True), _s(1.0, 5.55, 11.3, 0.5, 15, "muted"), _s(1.0, 6.35, 11.3, 0.4, 13, "muted", False, "left", True)]},
            "assignment": {"content": [_s(1.0, 1.9, 6.0, 0.5, 15, "primary", True), _s(1.0, 2.5, 11.3, 1.0, 17, "text", True)] + _twocol(1.0, 6.8, 3.95, 5.5, 0.8, 2, 0.9, 17, "text")},
        },
    },

    # ---------- 5. 中国文化：纸墨竖排，红竖线 ----------
    "lessonforge_deck_chinese_culture": {
        "cover": {
            "bg": "background",
            "deco": [("rect", 0, 0, 0.3, 7.5, "primary"), ("rect", 0.3, 0, 0.08, 7.5, "secondary"), ("rect", 1.6, 4.0, 2.4, 0.09, "primary")],
            "brand": (1.6, 6.9, 4, 0.4, "muted"),
            "title": _s(1.6, 2.35, 10.9, 1.6, 40, "primary", True, "center"),
            "subtitle": _s(1.6, 4.15, 10.9, 0.7, 18, "muted", False, "center"),
        },
        "end": {
            "bg": "primary",
            "deco": [("rect", 0, 0, 0.3, 7.5, "secondary"), ("oval", 11.3, -1.2, 3.2, 3.2, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(1.2, 2.4, 11.0, 1.3, 34, "on_primary", True, "center"),
            "subtitle": _s(1.2, 3.95, 11.0, 0.8, 20, "secondary", False, "center"),
        },
        "shell": {
            "deco": [("rect", 0, 0, 0.18, 7.5, "primary")],
            "brand": (0.5, 7.02, 4, 0.3, "muted"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(0.6, 0.55, 12.1, 0.85, 30, "primary", True, "center"),
            "accent": ("rect", 5.9, 1.42, 1.5, 0.09, "primary"),
        },
        "roles": {
            "intro": {"content": [_s(0.7, 1.85, 12.0, 1.0, 18, "primary", True, "center")] + _stack(0.7, 3.1, 12.0, 0.7, 5, 0.6, 17, "text", False, "center")},
            "objectives": {"content": [_s(0.7, 1.9, 12.0, 0.7, 18, "primary", True, "center")] + _stack(0.7, 2.8, 12.0, 0.85, 4, 0.75, 17, "text", False, "center")},
            "knowledge_map": {"content": _stack(0.7, 2.0, 12.0, 0.85, 5, 0.75, 17, "text", False, "center")},
            "knowledge_intro": {"content": [_s(0.7, 1.85, 12.0, 1.0, 18, "primary", True, "center")] + _stack(0.7, 3.1, 12.0, 0.62, 6, 0.55, 17, "text", False, "center")},
            "core_1": {"content": [_s(0.7, 1.85, 12.0, 0.5, 15, "primary", True, "center"), _s(0.7, 2.45, 12.0, 0.8, 18, "text", True, "center")] + _stack(0.7, 3.45, 12.0, 0.62, 4, 0.7, 17, "text", False, "center")},
            "core_2": {"content": [_s(0.7, 1.85, 12.0, 0.5, 15, "primary", True, "center"), _s(0.7, 2.45, 12.0, 0.8, 18, "text", True, "center")] + _stack(0.7, 3.45, 12.0, 0.62, 4, 0.7, 17, "text", False, "center")},
            "core_3": {"content": [_s(0.7, 1.85, 12.0, 0.5, 15, "primary", True, "center"), _s(0.7, 2.45, 12.0, 0.8, 18, "text", True, "center")] + _stack(0.7, 3.45, 12.0, 0.62, 4, 0.7, 17, "text", False, "center")},
            "core_4": {"content": [_s(0.7, 1.85, 12.0, 0.5, 15, "primary", True, "center"), _s(0.7, 2.45, 12.0, 0.8, 18, "text", True, "center")] + _stack(0.7, 3.45, 12.0, 0.62, 4, 0.7, 17, "text", False, "center")},
            "case_study": {"content": _grid(0.7, 1.9, 2.9, 1.1, 4, 1, 3.1, 0, 17, "text", True, "center") + [_s(0.7, 3.4, 12.0, 0.5, 15, "primary", True, "center"), _s(0.7, 4.0, 12.0, 1.2, 17, "text", False, "center")]},
            "discussion": {"content": [_s(0.7, 2.0, 12.0, 0.9, 18, "primary", True, "center")] + _stack(0.7, 3.3, 12.0, 0.8, 3, 0.95, 17, "text", False, "center")},
            "summary": {"content": _stack(0.7, 1.95, 12.0, 0.62, 8, 0.5, 17, "text", False, "center") + [_s(0.7, 6.1, 12.0, 0.6, 15, "muted", False, "center", True)]},
            "assessment": {"content": [_s(1.0, 1.9, 11.3, 0.5, 16, "text", True, "center"), _s(1.0, 2.5, 11.3, 0.5, 15, "muted", False, "center"), _s(1.0, 3.4, 11.3, 0.5, 16, "text", True, "center"), _s(1.0, 4.0, 11.3, 0.5, 15, "muted", False, "center"), _s(1.0, 4.9, 11.3, 0.5, 16, "text", True, "center"), _s(1.0, 5.5, 11.3, 0.5, 15, "muted", False, "center"), _s(1.0, 6.3, 11.3, 0.4, 13, "muted", False, "center", True)]},
            "assignment": {"content": [_s(1.0, 1.9, 11.3, 0.5, 15, "primary", True, "center"), _s(1.0, 2.55, 11.3, 1.0, 17, "text", False, "center")] + _stack(1.0, 3.85, 11.3, 0.65, 4, 0.55, 17, "text", False, "center")},
        },
    },

    # ---------- 6. 智慧课堂：紫色渐变侧栏 + 分栏面板 ----------
    "lessonforge_deck_smart_ai": {
        "cover": {
            "bg": "primary",
            "deco": [("rect", 0, 0, 2.6, 7.5, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(3.2, 2.5, 9.4, 1.7, 42, "on_primary", True),
            "subtitle": _s(3.2, 4.4, 9.4, 0.7, 20, "secondary"),
        },
        "end": {
            "bg": "primary",
            "deco": [("rect", 0, 0, 2.6, 7.5, "secondary"), ("oval", 11.4, -1.2, 3.4, 3.4, "secondary")],
            "brand": (0.95, 6.9, 4, 0.4, "on_primary"),
            "title": _s(3.2, 2.4, 9.4, 1.3, 34, "on_primary", True),
            "subtitle": _s(3.2, 3.95, 9.4, 0.8, 20, "secondary"),
        },
        "shell": {
            "deco": [("rect", 0, 0, 1.7, 7.5, "primary"), ("rect", 1.7, 0, 0.07, 7.5, "secondary")],
            "brand": (0.45, 7.02, 1.2, 0.3, "on_primary"),
            "page": (11.6, 7.02, 1.2, 0.3, "muted"),
            "title": _s(2.0, 0.55, 10.9, 0.85, 28, "primary", True),
            "accent": ("rect", 2.03, 1.46, 1.35, 0.09, "primary"),
        },
        "roles": {
            "intro": {"content": [_s(2.0, 1.85, 10.9, 1.0, 18, "primary", True)] + _twocol(2.0, 7.8, 3.1, 5.0, 0.95, 3, 0.5, 16, "text")},
            "objectives": {"content": _twocol(2.0, 7.8, 1.95, 5.0, 1.2, 3, 0.5, 16, "text")},
            "knowledge_map": {"content": _twocol(2.0, 7.8, 1.95, 5.0, 1.2, 3, 0.5, 16, "text")},
            "knowledge_intro": {"content": [_s(2.0, 1.85, 10.9, 0.95, 18, "primary", True)] + _twocol(2.0, 7.8, 3.05, 5.0, 0.8, 3, 0.5, 16, "text")},
            "core_1": {"content": [_s(2.0, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(2.0, 7.8, 2.5, 5.0, 1.1, 3, 0.5, 16, "text")},
            "core_2": {"content": [_s(2.0, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(2.0, 7.8, 2.5, 5.0, 1.1, 3, 0.5, 16, "text")},
            "core_3": {"content": [_s(2.0, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(2.0, 7.8, 2.5, 5.0, 1.1, 3, 0.5, 16, "text")},
            "core_4": {"content": [_s(2.0, 1.85, 3.2, 0.5, 14, "primary", True)] + _twocol(2.0, 7.8, 2.5, 5.0, 1.1, 3, 0.5, 16, "text")},
            "case_study": {"content": _grid(2.0, 1.85, 2.7, 1.1, 4, 1, 2.9, 0, 16, "text", True) + [_s(2.0, 3.35, 3.0, 0.5, 14, "primary", True), _s(2.0, 3.95, 10.9, 1.2, 16, "text")]},
            "discussion": {"content": [_s(2.0, 1.95, 10.9, 0.85, 18, "primary", True)] + _twocol(2.0, 7.8, 3.1, 5.0, 0.95, 2, 0.6, 16, "text")},
            "summary": {"content": [_s(2.0, 1.8, 10.9, 0.6, 17, "primary", True)] + _twocol(2.0, 7.8, 2.6, 5.0, 1.0, 3, 0.5, 16, "text") + [_s(2.0, 6.0, 10.9, 0.7, 15, "muted", False, "left", True)]},
            "assessment": {"content": [_s(2.0, 1.85, 10.9, 0.5, 15, "text", True), _s(2.0, 2.4, 10.9, 0.45, 14, "muted"), _s(2.0, 3.3, 10.9, 0.5, 15, "text", True), _s(2.0, 3.85, 10.9, 0.45, 14, "muted"), _s(2.0, 4.75, 10.9, 0.5, 15, "text", True), _s(2.0, 5.3, 10.9, 0.45, 14, "muted"), _s(2.0, 6.2, 10.9, 0.4, 12, "muted", False, "left", True)]},
            "assignment": {"content": [_s(2.0, 1.9, 6.0, 0.5, 14, "primary", True), _s(2.0, 2.5, 10.9, 1.0, 17, "text", True)] + _twocol(2.0, 7.8, 3.95, 5.0, 0.72, 2, 0.9, 16, "text")},
        },
    },
}


def hex_color(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def add_shape(slide, kind, x, y, w, h, color):
    if kind == "rect":
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, fill, border):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, spec, colors, fonts):
    x, y, w, h, size, color_key, bold, align, italic = spec
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    }[align]
    run = paragraph.add_run()
    run.text = ""
    run.font.name = fonts["body"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colors[color_key]
    return box


def _style_run(shape, color, fonts, size, bold, align):
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = fonts["body"]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_brand_page(prs, blank_layout, spec, colors, fonts, index):
    """整页背景页（封面/结课）：背景 + 装饰 + 品牌 + 标题 + 副标题。"""
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors[spec["bg"]]
    for kind, x, y, w, h, color_key in spec.get("deco", []):
        add_shape(slide, kind, x, y, w, h, colors[color_key])
    brand_x, brand_y, brand_w, brand_h, brand_color = spec["brand"]
    brand = slide.shapes.add_textbox(Inches(brand_x), Inches(brand_y), Inches(brand_w), Inches(brand_h))
    brand.text_frame.paragraphs[0].add_run().text = "LESSONFORGE 微课"
    _style_run(brand, colors[brand_color], fonts, 10, False, PP_ALIGN.LEFT)
    add_textbox(slide, spec["title"], colors, fonts)
    add_textbox(slide, spec["subtitle"], colors, fonts)
    return slide


def add_content_page(prs, blank_layout, design, role, spec, colors, fonts, index):
    """内容页：装饰 + 品牌 + 页码 + 标题 + 强调条 + 内容槽位（可选圆角卡片）。"""
    slide = prs.slides.add_slide(blank_layout)
    shell = design["shell"]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colors["background"]
    for kind, x, y, w, h, color_key in shell.get("deco", []):
        add_shape(slide, kind, x, y, w, h, colors[color_key])
    brand_x, brand_y, brand_w, brand_h, brand_color = shell["brand"]
    brand = slide.shapes.add_textbox(Inches(brand_x), Inches(brand_y), Inches(brand_w), Inches(brand_h))
    brand.text_frame.paragraphs[0].add_run().text = "LESSONFORGE 微课"
    _style_run(brand, colors[brand_color], fonts, 9, False, PP_ALIGN.LEFT)
    page_x, page_y, page_w, page_h, page_color = shell["page"]
    page = slide.shapes.add_textbox(Inches(page_x), Inches(page_y), Inches(page_w), Inches(page_h))
    page.text_frame.paragraphs[0].add_run().text = f"{index + 1:02d}"
    _style_run(page, colors[page_color], fonts, 10, True, PP_ALIGN.RIGHT)
    add_textbox(slide, shell["title"], colors, fonts)
    accent = shell.get("accent")
    if accent:
        kind, x, y, w, h, color_key = accent
        add_shape(slide, kind, x, y, w, h, colors[color_key])
    card = shell.get("card")
    for content in spec["content"]:
        if card:
            x, y, w, h = content[0], content[1], content[2], content[3]
            add_rounded_rect(
                slide, x - card["pad"], y - card["pad"],
                w + card["pad"] * 2, h + card["pad"] * 2,
                colors[card["fill"]], colors[card["border"]],
            )
        add_textbox(slide, content, colors, fonts)
    return slide


def build_deck(entry: dict) -> None:
    palette = entry["palette"]
    fonts = entry["typography"]
    colors = {key: hex_color(value) for key, value in palette.items()}
    deck_name = Path(entry["file"]).name
    design = TEMPLATE_DESIGNS[entry["id"]]

    prs = Presentation()
    prs.slide_width = Emu(int(CANVAS_W * EMU_PER_INCH))
    prs.slide_height = Emu(int(CANVAS_H * EMU_PER_INCH))
    blank_layout = prs.slide_layouts[6]
    for index, role in enumerate(ROLE_ORDER):
        if role == "cover":
            add_brand_page(prs, blank_layout, design["cover"], colors, fonts, index)
        elif role == "end":
            add_brand_page(prs, blank_layout, design["end"], colors, fonts, index)
        else:
            add_content_page(prs, blank_layout, design, role, design["roles"][role], colors, fonts, index)
    out = DECK_DIR / deck_name
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"  rebuilt {deck_name}")


def write_deck_slots() -> None:
    templates = {}
    for template_id, design in TEMPLATE_DESIGNS.items():
        roles = {"cover": {
            "title": {"x": design["cover"]["title"][0], "y": design["cover"]["title"][1]},
            "content": [{"x": design["cover"]["subtitle"][0], "y": design["cover"]["subtitle"][1]}],
        }}
        for role in ROLE_ORDER[1:-1]:
            spec = design["roles"][role]
            roles[role] = {
                "title": {"x": design["shell"]["title"][0], "y": design["shell"]["title"][1]},
                "content": [{"x": item[0], "y": item[1]} for item in spec["content"]],
            }
        roles["end"] = {
            "title": {"x": design["end"]["title"][0], "y": design["end"]["title"][1]},
            "content": [{"x": design["end"]["subtitle"][0], "y": design["end"]["subtitle"][1]}],
        }
        templates[template_id] = roles
    payload = {
        "version": "3.0.0",
        "role_order": ROLE_ORDER,
        "default_tol": 0.35,
        "templates": templates,
    }
    SLOTS_PATH.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"  wrote {SLOTS_PATH.name} (templates={len(templates)})")


def main() -> None:
    catalog = json.loads(CATALOG_PATH.read_text(encoding="utf-8"))
    print(f"rebuilding 6 distinct-design decks from {CATALOG_PATH.name}")
    for entry in catalog["templates"]:
        if entry["id"] not in TEMPLATE_DESIGNS:
            raise RuntimeError(f"catalog 模板 {entry['id']} 缺少设计系统定义")
        build_deck(entry)
    write_deck_slots()


if __name__ == "__main__":
    main()
